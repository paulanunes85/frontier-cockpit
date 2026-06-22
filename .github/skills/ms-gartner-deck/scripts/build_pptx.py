#!/usr/bin/env python3
"""
build_pptx.py: MS Gartner Deck -> native, fully-editable PowerPoint.

Every element is a native PowerPoint shape: text frames, native tables,
autoshapes, connectors. No images, nothing flattened. Open the .pptx and
every word, color, number, and box is editable.

The script does NOT parse HTML. It carries its own data model: a DECK list
of slide spec dicts. Each spec has a `type` mapped to one of 30 archetype
renderers. Edit DECK to change content; edit a renderer to change layout.

Usage:
    python3 build_pptx.py --out deck.pptx
    python3 build_pptx.py --out deck.pptx --only cover,table,quadrant
"""

import argparse
import sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Emu, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
    from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
    from pptx.oxml.ns import qn
    from lxml import etree
except ImportError as exc:
    sys.exit(
        f"ERROR: required dependency missing ({exc.name}).\n"
        "  pip install -r scripts/requirements.txt\n"
        "  (this script needs python-pptx and lxml)"
    )

# ============================================================
# CANVAS: 16:9 at 13.333" x 7.5"
# ============================================================
EMU_IN = 914400
def IN(v): return Emu(int(v * EMU_IN))
PW, PH = 13.333, 7.5

MARGIN_X = 0.62
TITLE_TOP = 0.86
CONTENT_TOP = 1.92
CONTENT_BOT = 6.74
CONTENT_W = PW - 2 * MARGIN_X
CONTENT_H = CONTENT_BOT - CONTENT_TOP

# ============================================================
# DESIGN TOKENS: mirror the HTML :root, Microsoft palette
# ============================================================
INK      = RGBColor(0x1A,0x1A,0x1A)
GRAPHITE = RGBColor(0x33,0x33,0x33)
SLATE    = RGBColor(0x5A,0x5A,0x5A)
SILVER   = RGBColor(0x8A,0x8A,0x8A)
PEARL    = RGBColor(0xB8,0xB8,0xB8)
LINE     = RGBColor(0xE4,0xE4,0xE4)
LINE2    = RGBColor(0xEF,0xEF,0xEF)
WHITE    = RGBColor(0xFF,0xFF,0xFF)
BG_SOFT  = RGBColor(0xF7,0xF7,0xF7)
BG_CARD  = RGBColor(0xFA,0xFA,0xFA)
BG_TINT  = RGBColor(0xF1,0xF1,0xF1)

MS = {
    "blue":   {"base": RGBColor(0x00,0xA4,0xEF), "d700": RGBColor(0x00,0x76,0xAC), "l50": RGBColor(0xE4,0xF6,0xFD)},
    "green":  {"base": RGBColor(0x7F,0xBA,0x00), "d700": RGBColor(0x5A,0x85,0x00), "l50": RGBColor(0xF1,0xF7,0xE1)},
    "yellow": {"base": RGBColor(0xFF,0xB9,0x00), "d700": RGBColor(0xB8,0x85,0x00), "l50": RGBColor(0xFF,0xF6,0xE0)},
    "red":    {"base": RGBColor(0xF2,0x50,0x22), "d700": RGBColor(0xB3,0x38,0x16), "l50": RGBColor(0xFD,0xEE,0xE9)},
}
LOGO_ORDER = ["red", "green", "blue", "yellow"]   # 2x2 logo squares

SECTION_COLOR = {
    "front": "blue", "ctx": "blue", "prob": "green", "sol": "yellow",
    "arch": "red", "plan": "blue", "invest": "green", "back": "red",
}
SECTION_LABEL = {
    "front": "Front matter", "ctx": "Section 1 \u00b7 Context",
    "prob": "Section 2 \u00b7 The problem", "sol": "Section 3 \u00b7 The solution",
    "arch": "Section 4 \u00b7 Architecture", "plan": "Section 5 \u00b7 The plan",
    "invest": "Section 6 \u00b7 Investment", "back": "Appendix",
}

SANS = "Segoe UI"     # Microsoft house font; always present in Office
MONO = "Consolas"

# ============================================================
# LOW-LEVEL SHAPE HELPERS
# ============================================================
def _fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color

def _noline(shape):
    shape.line.fill.background()

def _line(shape, color, w=1.0):
    shape.line.color.rgb = color
    shape.line.width = Pt(w)

def _noshadow(shape):
    spPr = shape._element.spPr
    if spPr.find(qn('a:effectLst')) is None:
        etree.SubElement(spPr, qn('a:effectLst'))

def rect(slide, x, y, w, h, fill=None, line=None, line_w=1.0,
         rounded=False, radius=0.055):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
        IN(x), IN(y), IN(w), IN(h))
    if rounded:
        try: shp.adjustments[0] = radius
        except Exception: pass
    if fill is None: shp.fill.background()
    else: _fill(shp, fill)
    if line is None: _noline(shp)
    else: _line(shp, line, line_w)
    _noshadow(shp)
    return shp

def oval(slide, x, y, w, h, fill, line=None, line_w=1.0):
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL, IN(x), IN(y), IN(w), IN(h))
    _fill(shp, fill)
    if line is None: _noline(shp)
    else: _line(shp, line, line_w)
    _noshadow(shp)
    return shp

def hline(slide, x, y, w, color, weight=1.0):
    ln = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, IN(x), IN(y), IN(x+w), IN(y))
    ln.line.color.rgb = color
    ln.line.width = Pt(weight)
    return ln

def vline(slide, x, y, h, color, weight=1.0):
    ln = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, IN(x), IN(y), IN(x), IN(y+h))
    ln.line.color.rgb = color
    ln.line.width = Pt(weight)
    return ln

def text(slide, x, y, w, h, runs, *, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         line_spacing=1.0, wrap=True):
    """runs: string, or list of paragraphs. Each paragraph is a list of
    (text, opts) tuples. opts: size, bold, italic, color, font, tracking
    (letter-spacing pt), caps, align, space_before, space_after, line_spacing."""
    box = slide.shapes.add_textbox(IN(x), IN(y), IN(w), IN(h))
    tf = box.text_frame
    tf.word_wrap = wrap
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    if isinstance(runs, str):
        runs = [[(runs, {})]]
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        pa = para[0][1] if para and len(para[0]) > 1 else {}
        p.line_spacing = pa.get("line_spacing", line_spacing)
        if pa.get("align"): p.alignment = pa["align"]
        if pa.get("space_before"): p.space_before = Pt(pa["space_before"])
        if pa.get("space_after"): p.space_after = Pt(pa["space_after"])
        for seg in para:
            t, o = (seg[0], seg[1]) if len(seg) > 1 else (seg[0], {})
            r = p.add_run()
            r.text = t.upper() if o.get("caps") else t
            f = r.font
            f.size = Pt(o.get("size", 12))
            f.bold = o.get("bold", False)
            f.italic = o.get("italic", False)
            f.name = o.get("font", SANS)
            f.color.rgb = o.get("color", INK)
            tr = o.get("tracking")
            if tr is not None:
                r._r.get_or_add_rPr().set("spc", str(int(tr * 100)))
    return box

def native_table(slide, x, y, w, rows, cols, col_widths=None, row_h=0.34):
    """Native PowerPoint table, default banding stripped. Returns table object."""
    gt = slide.shapes.add_table(rows, cols, IN(x), IN(y), IN(w), IN(row_h*rows))
    tbl = gt.table
    tblPr = tbl._tbl.find(qn('a:tblPr'))
    if tblPr is not None:
        tblPr.set('firstRow', '0')
        tblPr.set('bandRow', '0')
        for child in list(tblPr):
            tblPr.remove(child)
    if col_widths:
        for ci, cw in enumerate(col_widths):
            tbl.columns[ci].width = IN(cw)
    return tbl

def cell_text(cell, runs, *, size=11, bold=False, color=GRAPHITE,
              align=PP_ALIGN.LEFT, fill=None, font=SANS, anchor=MSO_ANCHOR.MIDDLE):
    """Fill one table cell. runs: string or list of (text, opts)."""
    if fill is not None:
        cell.fill.solid()
        cell.fill.fore_color.rgb = fill
    else:
        cell.fill.background()
    cell.vertical_anchor = anchor
    cell.margin_left = IN(0.10); cell.margin_right = IN(0.10)
    cell.margin_top = IN(0.04); cell.margin_bottom = IN(0.04)
    tf = cell.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    if isinstance(runs, str):
        runs = [(runs, {})]
    for seg in runs:
        if isinstance(seg, tuple):
            t, o = (seg[0], seg[1]) if len(seg) > 1 else (seg[0], {})
        else:
            t, o = seg, {}
        r = p.add_run()
        r.text = t
        f = r.font
        f.size = Pt(o.get("size", size))
        f.bold = o.get("bold", bold)
        f.name = o.get("font", font)
        f.color.rgb = o.get("color", color)


# ============================================================
# CHROME: logo, top bar, title block, footer, base slide
# ============================================================
def ms_logo(slide, x, y, size=0.22):
    """2x2 Microsoft squares as 4 native rectangles."""
    gap = size * 0.09
    cell = (size - gap) / 2
    for (cx, cy), key in zip([(0,0),(1,0),(0,1),(1,1)], LOGO_ORDER):
        rect(slide, x + cx*(cell+gap), y + cy*(cell+gap), cell, cell,
             fill=MS[key]["base"])

def chrome_top(slide, act_label, page_no, total):
    ms_logo(slide, MARGIN_X, 0.40, size=0.22)
    text(slide, MARGIN_X + 0.34, 0.40, 7.0, 0.24, [[
        ("Microsoft", {"size": 9.5, "bold": True, "color": INK, "tracking": 1.4, "caps": True}),
        ("    \u00b7    ", {"size": 9.5, "color": PEARL, "tracking": 1.4}),
        (act_label, {"size": 9.5, "color": SLATE, "tracking": 1.4, "caps": True}),
    ]], anchor=MSO_ANCHOR.MIDDLE)
    text(slide, PW - MARGIN_X - 2.0, 0.40, 2.0, 0.24, [[
        (f"{page_no:02d}", {"size": 10, "bold": True, "color": INK}),
        ("  \u00b7  ", {"size": 10, "color": PEARL}),
        (f"{total}", {"size": 10, "color": SLATE}),
    ]], align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
    hline(slide, MARGIN_X, 0.74, CONTENT_W, LINE, 0.75)

def title_block(slide, kicker, title_runs, descriptor, accent, title_size=22):
    """Kicker + title + descriptor. Title height adapts to estimated wrap.
    Returns the y where content can safely begin."""
    text(slide, MARGIN_X, TITLE_TOP, CONTENT_W, 0.22, [[
        (kicker, {"size": 9.5, "bold": True, "color": accent["d700"],
                  "tracking": 2.2, "caps": True}),
    ]])
    # estimate title line count from weighted character length
    # (bold runs render wider, weight them up)
    title_chars = 0
    for seg in title_runs:
        w = 1.18 if seg[1].get("bold") else 1.0
        title_chars += len(seg[0]) * w
    # at title_size, ~ this many weighted chars fit on one CONTENT_W line
    chars_per_line = (CONTENT_W * 132.0) / title_size
    import math
    title_lines = max(1, math.ceil(title_chars / chars_per_line))
    line_h = title_size * 1.18 / 72.0   # inches per line
    title_h = title_lines * line_h + 0.08
    text(slide, MARGIN_X, TITLE_TOP + 0.24, CONTENT_W, title_h + 0.2,
         [title_runs], line_spacing=1.14)
    desc_y = TITLE_TOP + 0.24 + title_h + 0.08
    if descriptor:
        text(slide, MARGIN_X, desc_y, CONTENT_W - 0.8, 0.5,
             [descriptor], line_spacing=1.3)
        desc_y += 0.46
    return desc_y

def chrome_bottom(slide, src_runs):
    hline(slide, MARGIN_X, CONTENT_BOT + 0.12, CONTENT_W, LINE, 0.75)
    text(slide, MARGIN_X, CONTENT_BOT + 0.20, CONTENT_W - 2.7, 0.32,
         [src_runs], line_spacing=1.25)
    text(slide, PW - MARGIN_X - 2.6, CONTENT_BOT + 0.20, 2.6, 0.32, [[
        ("Microsoft \u00b7 Strategic Brief", {"size": 8, "bold": True,
         "color": SLATE, "tracking": 1.5, "caps": True}),
    ]], align=PP_ALIGN.RIGHT)

def base_slide(prs, bg=WHITE):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    b = rect(slide, 0, 0, PW, PH, fill=bg)
    spTree = slide.shapes._spTree
    spTree.remove(b._element)
    spTree.insert(2, b._element)
    return slide

def content_frame(prs, sec, kicker, title_runs, descriptor, src_runs, title_size=22):
    """Standard content slide: chrome + title + footer.
    Returns (slide, accent, content_top), content_top adapts to title length."""
    slide = base_slide(prs)
    accent = MS[SECTION_COLOR[sec]]
    chrome_top(slide, SECTION_LABEL[sec], CUR["page"], CUR["total"])
    ctop = title_block(slide, kicker, title_runs, descriptor, accent, title_size)
    chrome_bottom(slide, src_runs)
    # ensure a minimum gap below the title; allow 2-line titles to push down
    ctop = max(ctop + 0.10, CONTENT_TOP)
    return slide, accent, ctop

def bleed_slide(prs, bg=WHITE):
    """Full-bleed slide, no chrome."""
    return base_slide(prs, bg)

# CUR holds the running page/total so renderers don't each take them
CUR = {"page": 1, "total": 53}

# inline-run helper: turns "plain **bold** plain" markup into run tuples
def runs(s, size=12, color=INK, bold_color=None, **base):
    """Parse **bold** and *mark* markup into a list of (text, opts) tuples."""
    bold_color = bold_color or color
    out = []
    i = 0
    while i < len(s):
        if s[i:i+2] == "**":
            j = s.index("**", i+2)
            out.append((s[i+2:j], {"size": size, "bold": True, "color": bold_color, **base}))
            i = j + 2
        else:
            nxt = s.find("**", i)
            seg = s[i:] if nxt == -1 else s[i:nxt]
            if seg:
                out.append((seg, {"size": size, "color": color, **base}))
            i = len(s) if nxt == -1 else nxt
    return out


def note_plain(value):
    """Plain text for PowerPoint speaker notes."""
    text_value = str(value or "")
    text_value = text_value.replace("**", "").replace("*", "")
    text_value = text_value.replace("\u00b7", " ")
    return " ".join(text_value.split())


def collect_note_text(value):
    """Collect meaningful strings from a slide spec for note synthesis."""
    if isinstance(value, str):
        text_value = note_plain(value)
        return [text_value] if len(text_value) > 12 else []
    if isinstance(value, (list, tuple)):
        items = []
        for item in value:
            items.extend(collect_note_text(item))
        return items
    if isinstance(value, dict):
        skip = {"type", "sec", "q_start", "act_no"}
        items = []
        for key, item in value.items():
            if key not in skip:
                items.extend(collect_note_text(item))
        return items
    return []


def speaker_note(spec, page, total):
    """Generate an English presenter note from the slide spec."""
    title = note_plain(spec.get("h") or spec.get("title") or spec.get("quote") or spec.get("statement") or spec.get("type", "Slide"))
    snippets = []
    for item in collect_note_text(spec):
        if item != title and item not in snippets:
            snippets.append(item)
    core = " ".join(snippets[:3]) or title
    if len(core) > 520:
        core = core[:517].rstrip() + "."
    return (
        f"[OPENING]\nSlide {page} of {total}: {title}\n\n"
        f"[CORE]\nUse this slide to advance the strategic argument. Ground the room in the visible structure, then explain: {core}\n\n"
        "[PROVOCATIVE HOOK]\nAsk what decision the audience can make with this evidence that it could not make before this slide.\n\n"
        "[TRANSITION]\nClose by naming the next step in the story so the deck feels like one argument, not a set of disconnected pages.\n\n"
        "[TIMING] ~1 min"
    )


def attach_speaker_note(slide, spec, page, total):
    slide.notes_slide.notes_text_frame.text = speaker_note(spec, page, total)


# ============================================================
# RENDERERS: cover (4 variants)
# Left panel is fixed; right panel swaps by variant.
# ============================================================
def _cover_left(slide):
    half = PW * 0.52
    rect(slide, 0, 0, half, PH, fill=INK)
    # 4-color top edge
    seg = half / 4
    for i, key in enumerate(["red","green","yellow","blue"]):
        rect(slide, i*seg, 0, seg, 0.07, fill=MS[key]["base"])
    # brand
    ms_logo(slide, 0.62, 0.54, size=0.24)
    text(slide, 0.98, 0.54, half-1.2, 0.28, [[
        ("Microsoft", {"size": 10, "bold": True, "color": WHITE, "tracking": 1.6, "caps": True}),
        ("   \u00b7   ", {"size": 10, "color": RGBColor(0x55,0x55,0x55)}),
        ("AI Platform Modernization", {"size": 10, "color": RGBColor(0xC8,0xC8,0xC8), "tracking": 1.6, "caps": True}),
    ]], anchor=MSO_ANCHOR.MIDDLE)
    # title, sized to wrap to 3 lines inside ~5.3in
    text(slide, 0.62, 1.30, half-0.92, 2.30, [[
        ("The platform decision behind every ", {"size": 27, "bold": True, "color": WHITE}),
        ("production AI", {"size": 27, "bold": True, "color": MS["blue"]["base"]}),
        (" outcome.", {"size": 27, "bold": True, "color": WHITE}),
    ]], line_spacing=1.08)
    # subtitle
    text(slide, 0.62, 3.70, half-1.15, 1.0, [
        runs("A strategic and technical brief on why enterprise AI initiatives "
             "stall before production, and the reference architecture that "
             "closes the gap between pilot and scale.", size=12.5,
             color=RGBColor(0xC8,0xC8,0xC8))
    ], line_spacing=1.4)
    # chips, single row, sized to fit
    chips = [("Azure ", "AI Foundry"), ("GitHub ", "Copilot Enterprise"),
             ("Azure ", "Kubernetes Service"), ("Agentic ", "DevOps")]
    cx, cy = 0.62, 4.78
    for pre, strong in chips:
        label = pre + strong
        cw = 0.24 + len(label) * 0.066
        rect(slide, cx, cy, cw, 0.33, fill=None, line=RGBColor(0x44,0x44,0x44),
             line_w=1.0, rounded=True, radius=0.5)
        text(slide, cx, cy, cw, 0.33, [[
            (pre, {"size": 10, "color": RGBColor(0xE0,0xE0,0xE0)}),
            (strong, {"size": 10, "bold": True, "color": WHITE}),
        ]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        cx += cw + 0.11
    # byline
    hline(slide, 0.62, 6.34, half-1.2, RGBColor(0x3A,0x3A,0x3A), 0.75)
    bw = (half - 1.2) / 2
    for i, (name, role) in enumerate([
        ("Frontier Cockpit Team", "Software Global Black Belt"),
        ("Strategic & Technical Brief", "Client engagement \u00b7 2026")]):
        text(slide, 0.62 + i*bw, 6.50, bw, 0.7, [
            [(name, {"size": 11.5, "bold": True, "color": WHITE})],
            [(role, {"size": 8.5, "color": RGBColor(0x9A,0x9A,0x9A), "tracking": 1.1,
              "caps": True, "space_before": 3})],
        ], line_spacing=1.1)
    return half

def _cover_right_head(slide, half, top_label):
    rx = half + 0.55
    rw = PW - half - 1.15
    text(slide, rx, 0.60, rw, 0.26, [[
        (top_label, {"size": 9.5, "bold": True, "color": SLATE, "tracking": 1.6, "caps": True}),
    ]])
    text(slide, rx, 0.60, rw, 0.26, [[
        (f"{CUR['page']:02d} / {CUR['total']}", {"size": 9.5, "bold": True,
         "color": MS["blue"]["base"]}),
    ]], align=PP_ALIGN.RIGHT)
    return rx, rw

def r_cover_meeting(prs, s):
    slide = base_slide(prs)
    half = _cover_left(slide)
    rx, rw = _cover_right_head(slide, half, "Confidential \u00b7 Client brief")
    text(slide, rx, 1.42, rw, 0.22, [[
        ("Session context", {"size": 10, "bold": True, "color": MS["blue"]["d700"],
         "tracking": 2.4, "caps": True})]])
    text(slide, rx, 1.66, rw, 1.10, [
        runs("A working session to decide the **platform scope and the "
             "first phase**.", size=18, color=INK, bold_color=INK)
    ], line_spacing=1.24)
    # 2x2 meeting grid
    grid = [("Client", "Contoso Financial", "Platform & Engineering leadership"),
            ("Date", "2026-05-20", "90-minute briefing"),
            ("Session type", "Strategic & technical brief", "Decision-oriented, not a pitch"),
            ("In the room", "CTO, Head of Platform", "Security and SRE leads")]
    hline(slide, rx, 2.92, rw, LINE, 0.75)
    gw = rw / 2
    gh = 0.92
    for i, (k, v, sub) in enumerate(grid):
        gx = rx + (i % 2) * gw
        gy = 3.06 + (i // 2) * gh
        text(slide, gx, gy, gw - 0.24, gh, [
            [(k, {"size": 9, "bold": True, "color": SLATE, "tracking": 1.5, "caps": True})],
            [(v, {"size": 13.5, "bold": True, "color": INK, "space_before": 4})],
            [(sub, {"size": 10.5, "color": SLATE, "space_before": 2})],
        ], line_spacing=1.18)
    # outcome bar
    oy = 3.06 + 2*gh + 0.16
    rect(slide, rx, oy, 0.045, 0.74, fill=MS["blue"]["base"])
    text(slide, rx + 0.18, oy, rw - 0.18, 0.74, [
        runs("**Outcome we are after:** agreement on a three-week assessment "
             "window and the two candidate workloads to baseline. Placeholders "
             "for the reference template.", size=11.5, color=SLATE, bold_color=GRAPHITE)
    ], line_spacing=1.42)

def r_cover_stat(prs, s):
    slide = base_slide(prs)
    half = _cover_left(slide)
    rx, rw = _cover_right_head(slide, half, "Confidential \u00b7 Client brief")
    text(slide, rx, 2.05, rw, 0.24, [[
        ("The opening number", {"size": 10, "bold": True, "color": MS["blue"]["d700"],
         "tracking": 2.4, "caps": True})]])
    text(slide, rx - 0.04, 2.20, rw, 1.6, [[
        ("3", {"size": 128, "bold": True, "color": INK}),
        ("\u00d7", {"size": 46, "color": SLATE}),
    ]], line_spacing=0.9)
    text(slide, rx, 3.95, rw - 0.2, 0.8, [
        runs("longer time-to-production for AI workloads built on **ad hoc "
             "infrastructure** versus an opinionated platform baseline.",
             size=15, color=INK, bold_color=INK)
    ], line_spacing=1.35)
    text(slide, rx, 4.85, rw - 0.2, 0.5, [
        runs("Illustrative figure for the reference template. Replace with "
             "**cited industry data** in real decks.", size=11, color=SLATE, bold_color=INK)
    ], line_spacing=1.3)
    # agenda row
    hline(slide, rx, 5.80, rw, LINE, 0.75)
    aw = rw / 3
    agenda = [("I.", "Context", "Where the market is", "blue"),
              ("II.", "Solution", "Platform as the unit of delivery", "green"),
              ("III.", "Execute", "Roadmap and investment", "red")]
    for i, (num, head, sub, ckey) in enumerate(agenda):
        ax = rx + i*aw
        text(slide, ax, 5.96, aw-0.16, 0.66, [
            [(num + " ", {"size": 12, "bold": True, "color": MS[ckey]["base"]}),
             (head, {"size": 12, "bold": True, "color": INK})],
            [(sub, {"size": 8.5, "color": SLATE, "tracking": 1.1, "caps": True, "space_before": 3})],
        ], line_spacing=1.15)

def r_cover_thesis(prs, s):
    slide = base_slide(prs)
    half = _cover_left(slide)
    rx, rw = _cover_right_head(slide, half, "Confidential \u00b7 Client brief")
    text(slide, rx, 1.45, rw, 0.22, [[
        ("The thesis", {"size": 10, "bold": True, "color": MS["blue"]["d700"],
         "tracking": 2.4, "caps": True})]])
    text(slide, rx, 1.74, rw, 2.9, [
        runs("The model is not the bottleneck. The **platform underneath it** "
             "is what decides whether a pilot ever becomes a service the "
             "business can depend on.", size=21, color=INK, bold_color=INK)
    ], line_spacing=1.26)
    rect(slide, rx, 4.95, 0.045, 1.05, fill=MS["blue"]["base"])
    text(slide, rx + 0.18, 4.95, rw - 0.18, 1.05, [
        runs("This brief makes the case for treating the platform as the unit "
             "of delivery, then lays out the architecture, roadmap, and "
             "investment to build one. **Content is illustrative** for the "
             "reference template.", size=12, color=SLATE, bold_color=GRAPHITE)
    ], line_spacing=1.45)

def r_cover_pillars(prs, s):
    slide = base_slide(prs)
    half = _cover_left(slide)
    rx, rw = _cover_right_head(slide, half, "Confidential \u00b7 Client brief")
    text(slide, rx, 1.95, rw, 0.24, [[
        ("Three horizons", {"size": 10, "bold": True, "color": MS["blue"]["d700"],
         "tracking": 2.4, "caps": True})]])
    text(slide, rx, 2.22, rw, 0.4, [[
        ("The brief is built on three moves, each earning the next.",
         {"size": 15, "color": INK})]], line_spacing=1.4)
    items = [("I", "Diagnose the gap", "Why pilots stall before production, sized with the client's own data.", "blue"),
             ("II", "Reframe the solution", "Platform as the unit of delivery, with an opinionated reference stack.", "green"),
             ("III", "Execute the plan", "A staged roadmap, a clear investment, and an honest view of the risks.", "red")]
    iy = 2.92
    for num, title, desc, ckey in items:
        rect(slide, rx, iy, rw, 0.92, fill=BG_CARD, line=LINE, line_w=1.0,
             rounded=True, radius=0.10)
        rect(slide, rx + 0.16, iy + 0.24, 0.44, 0.44, fill=MS[ckey]["base"],
             rounded=True, radius=0.20)
        text(slide, rx + 0.16, iy + 0.24, 0.44, 0.44, [[
            (num, {"size": 16, "bold": True, "color": WHITE if ckey != "yellow" else INK})
        ]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        text(slide, rx + 0.76, iy + 0.16, rw - 0.92, 0.7, [
            [(title, {"size": 14, "bold": True, "color": INK})],
            [(desc, {"size": 11, "color": SLATE, "space_before": 2})],
        ], line_spacing=1.32)
        iy += 1.04


# ============================================================
# RENDERERS: divider, big-number, data-table
# ============================================================
def r_divider(prs, s):
    slide = base_slide(prs)
    sec = s["sec"]
    accent = MS[SECTION_COLOR[sec]]
    half = PW / 2
    # left white panel
    rect(slide, 0, 0, half, PH, fill=WHITE)
    rect(slide, 0, 0, 0.08, PH, fill=accent["base"])  # accent spine
    text(slide, 0.62, 1.95, 3.0, 1.6, [[
        (s["act_no"], {"size": 112, "bold": True, "color": accent["base"]})
    ]], line_spacing=0.8)
    text(slide, 0.62, 3.50, half-1.0, 0.26, [[
        (s["kicker"], {"size": 10.5, "bold": True, "color": accent["d700"],
         "tracking": 2.4, "caps": True})]])
    text(slide, 0.62, 3.80, half-1.0, 1.0, [
        runs(s["title"], size=36, color=INK, bold_color=INK)
    ], line_spacing=1.12)
    text(slide, 0.62, 4.82, half-1.2, 0.9, [
        runs(s["desc"], size=13.5, color=SLATE)
    ], line_spacing=1.45)
    # right accent panel
    rect(slide, half, 0, half, PH, fill=accent["base"])
    on_yellow = SECTION_COLOR[sec] == "yellow"
    q_color = INK if on_yellow else WHITE
    mini_color = RGBColor(0x44,0x44,0x44) if on_yellow else RGBColor(0xE6,0xE6,0xE6)
    rule_color = RGBColor(0x00,0x00,0x00) if on_yellow else WHITE
    text(slide, half+0.7, 3.00, half-1.4, 0.24, [[
        ("Questions this act answers", {"size": 10, "bold": True,
         "color": mini_color, "tracking": 1.8, "caps": True})]])
    qy = 3.42
    for i, q in enumerate(s["questions"]):
        text(slide, half+0.7, qy, half-1.4, 0.42, [[
            (f"Q{s['q_start']+i}  ", {"size": 14, "bold": True, "color": q_color}),
            (q, {"size": 14, "color": q_color}),
        ]], line_spacing=1.3)
        if i < len(s["questions"]) - 1:
            ln = hline(slide, half+0.7, qy+0.46, half-1.4, rule_color, 0.5)
        qy += 0.62

def r_big_number(prs, s):
    slide, accent, ctop = content_frame(prs, s["sec"], s["kicker"],
        runs(s["h"], size=22, color=INK, bold_color=INK),
        runs(s["desc"], size=12.5, color=SLATE),
        runs(s["src"], size=10.5, color=SLATE, bold_color=GRAPHITE),
        title_size=22)
    # dark metric card (left)
    cy = ctop
    ch = CONTENT_BOT - cy - 0.05
    cw_l = CONTENT_W * 0.40
    rect(slide, MARGIN_X, cy, cw_l, ch, fill=GRAPHITE, rounded=True, radius=0.04)
    pad = 0.28
    text(slide, MARGIN_X+pad, cy+ch*0.20, cw_l-2*pad, 0.24, [[
        (s["metric_label"], {"size": 9.5, "bold": True,
         "color": RGBColor(0xC0,0xC0,0xC0), "tracking": 2.0, "caps": True})]])
    text(slide, MARGIN_X+pad, cy+ch*0.28, cw_l-2*pad, 1.1, [[
        (s["metric_n"], {"size": 64, "bold": True, "color": WHITE}),
        (s["metric_unit"], {"size": 28, "color": accent["base"]}),
    ]], line_spacing=0.9)
    text(slide, MARGIN_X+pad, cy+ch*0.62, cw_l-2*pad, ch*0.34, [
        runs(s["metric_desc"], size=12, color=RGBColor(0xDD,0xDD,0xDD), bold_color=WHITE)
    ], line_spacing=1.4)
    # 3 supporting cards (right)
    rx = MARGIN_X + cw_l + 0.32
    rw = CONTENT_W - cw_l - 0.32
    gap = 0.14
    sh = (ch - 2*gap) / 3
    for i, (lbl, n, d, klass) in enumerate(s["cards"]):
        sy = cy + i*(sh+gap)
        is_accent = klass == "accent"
        rect(slide, rx, sy, rw, sh, fill=accent["l50"] if is_accent else BG_CARD,
             line=accent["base"] if is_accent else LINE, line_w=1.0,
             rounded=True, radius=0.06)
        text(slide, rx+0.22, sy+0.16, rw-0.44, 0.20, [[
            (lbl, {"size": 9, "bold": True, "color": accent["d700"],
             "tracking": 1.8, "caps": True})]])
        text(slide, rx+0.22, sy+0.40, 1.5, 0.5, [[
            (n, {"size": 27, "bold": True, "color": INK})]])
        text(slide, rx+1.85, sy+0.30, rw-2.05, sh-0.5, [
            runs(d, size=11, color=SLATE, bold_color=INK)
        ], line_spacing=1.32, anchor=MSO_ANCHOR.MIDDLE)

# tag pill colors for tables
TAG_COLORS = {
    "High":    ("red",),    "Medium": ("yellow",),
    "Watch":   ("blue",),   "Low":    ("green",),
    "Partial": ("yellow",),
}
def _tag_pill(slide, x, y, label):
    """A small colored tag pill, sized to label, single-line."""
    key = TAG_COLORS.get(label, ("blue",))[0]
    w = 0.30 + len(label) * 0.078
    rect(slide, x, y, w, 0.26, fill=MS[key]["l50"], rounded=True, radius=0.5)
    box = text(slide, x, y, w, 0.26, [[
        (label, {"size": 8.5, "bold": True, "color": MS[key]["d700"],
         "tracking": 0.4, "caps": True})]],
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, wrap=False)
    return w

def _est_lines(s_text, col_w_in, font_pt):
    """Estimate wrapped line count. Calibrated against LibreOffice/PowerPoint,
    which render Segoe UI ~18% wider than a naive char count predicts."""
    if not s_text:
        return 1
    import math
    chars_per_in = 155.0 / font_pt    # conservative, renderer runs wide
    cap = max(1, int(col_w_in * chars_per_in))
    return max(1, math.ceil(len(s_text) / cap))

def r_data_table(prs, s):
    slide, accent, ctop = content_frame(prs, s["sec"], s["kicker"],
        runs(s["h"], size=22, color=INK, bold_color=INK),
        runs(s["desc"], size=12.5, color=SLATE),
        runs(s["src"], size=10.5, color=SLATE, bold_color=GRAPHITE),
        title_size=22)
    headers = s["headers"]   # (label, align, width_frac)
    rows = s["rows"]         # list of [(text, kind), ...]  kind: None|num|tag|bold
    ncols = len(headers)
    col_w = [h[2]*CONTENT_W for h in headers]
    col_x = [MARGIN_X]
    for cw in col_w[:-1]:
        col_x.append(col_x[-1] + cw)

    # --- estimate per-row heights from the tallest wrapping cell ---
    HEADER_H = 0.40
    LINE_H = 0.205          # height per wrapped text line
    PAD_V = 0.18            # vertical padding inside a row
    body_h = []
    for row in rows:
        max_lines = 1
        for ci, (val, kind) in enumerate(row):
            if kind in ("tag", "num"):
                continue
            fs = 11.5
            max_lines = max(max_lines, _est_lines(val, col_w[ci]-0.22, fs))
        body_h.append(max(0.40, max_lines*LINE_H + PAD_V))
    total_body = sum(body_h)
    ty = ctop

    # --- build native table with explicit row heights ---
    nrows = len(rows) + 1
    tbl = native_table(slide, MARGIN_X, ty, CONTENT_W, nrows, ncols,
                       col_widths=col_w, row_h=0.40)
    tbl.rows[0].height = IN(HEADER_H)
    for ri, rh in enumerate(body_h, start=1):
        tbl.rows[ri].height = IN(rh)
    # header
    for ci, (label, al, _) in enumerate(headers):
        cell_text(tbl.cell(0, ci), [(label, {})], size=9, bold=True, color=SLATE,
                  align=PP_ALIGN.RIGHT if al == "num" else PP_ALIGN.LEFT,
                  fill=WHITE, anchor=MSO_ANCHOR.BOTTOM)
    # body
    for ri, row in enumerate(rows, start=1):
        zebra = BG_SOFT if ri % 2 == 0 else WHITE
        for ci, (val, kind) in enumerate(row):
            c = tbl.cell(ri, ci)
            if kind == "tag":
                cell_text(c, "", fill=zebra)
            elif kind == "num":
                cell_text(c, [(val, {"bold": True})], size=11, color=GRAPHITE,
                          align=PP_ALIGN.RIGHT, fill=zebra)
            elif kind == "bold":
                cell_text(c, [(val, {"bold": True, "color": INK})], size=11.5, fill=zebra)
            else:
                cell_text(c, [(val, {})], size=11.5, color=GRAPHITE, fill=zebra)
    # heavy rule under header
    hline(slide, MARGIN_X, ty + HEADER_H, CONTENT_W, GRAPHITE, 1.5)
    # overlay tag pills at the true cumulative row positions
    y_cursor = ty + HEADER_H
    for ri, row in enumerate(rows):
        rh = body_h[ri]
        for ci, (val, kind) in enumerate(row):
            if kind == "tag":
                px = col_x[ci] + 0.10
                py = y_cursor + (rh - 0.24)/2
                _tag_pill(slide, px, py, val)
        y_cursor += rh
    table_bottom = ty + HEADER_H + total_body

    # optional pull quote, placed below the real table bottom
    if s.get("pull"):
        py = table_bottom + 0.20
        ph = CONTENT_BOT - py
        if ph > 0.70:
            lbl, txt, cite = s["pull"]
            rect(slide, MARGIN_X, py, CONTENT_W, ph, fill=accent["l50"],
                 rounded=True, radius=0.04)
            rect(slide, MARGIN_X, py, 0.045, ph, fill=accent["base"])
            text(slide, MARGIN_X+0.24, py+0.13, CONTENT_W-0.5, 0.18, [[
                (lbl, {"size": 9, "bold": True, "color": accent["d700"],
                 "tracking": 1.8, "caps": True})]])
            text(slide, MARGIN_X+0.24, py+0.34, CONTENT_W-0.6, ph-0.62, [
                runs(txt, size=12.5, color=INK, bold_color=INK)
            ], line_spacing=1.32)
            text(slide, MARGIN_X+0.24, py+ph-0.26, CONTENT_W-0.5, 0.20, [
                runs(cite, size=9.5, color=SLATE, bold_color=GRAPHITE)
            ])


# ============================================================
# RENDERERS: exec-summary, quote, section-intro, pillar-grid, flow
# ============================================================
def r_exec_summary(prs, s):
    slide, accent, ctop = content_frame(prs, s["sec"], s["kicker"],
        runs(s["h"], size=22, color=INK, bold_color=INK),
        runs(s["desc"], size=12.5, color=SLATE),
        runs(s["src"], size=10.5, color=SLATE, bold_color=GRAPHITE),
        title_size=22)
    cols = s["cols"]   # list of (number, title, body, take)
    n = len(cols)
    gap = 0.24
    cw = (CONTENT_W - gap*(n-1)) / n
    cy = ctop + 0.15
    ch = CONTENT_BOT - cy - 0.10
    num_colors = ["blue", "green", "red", "yellow"]
    for i, (num, title, body, take) in enumerate(cols):
        cx = MARGIN_X + i*(cw+gap)
        rect(slide, cx, cy, cw, ch, fill=BG_CARD, line=LINE, line_w=1.0,
             rounded=True, radius=0.045)
        pad = 0.26
        # number chip
        rect(slide, cx+pad, cy+pad, 0.34, 0.34, fill=MS[num_colors[i % 4]]["base"],
             rounded=True, radius=0.18)
        text(slide, cx+pad, cy+pad, 0.34, 0.34, [[
            (str(num), {"size": 13, "bold": True,
             "color": INK if num_colors[i%4]=="yellow" else WHITE})]],
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # title
        text(slide, cx+pad, cy+pad+0.48, cw-2*pad, 0.7, [
            runs(title, size=15, color=INK, bold_color=INK)
        ], line_spacing=1.2)
        # body
        text(slide, cx+pad, cy+pad+1.20, cw-2*pad, ch-pad-2.0, [
            runs(body, size=11.5, color=SLATE)
        ], line_spacing=1.45)
        # take-away pinned near bottom
        hline(slide, cx+pad, cy+ch-0.86, cw-2*pad, LINE2, 0.75)
        text(slide, cx+pad, cy+ch-0.74, cw-2*pad, 0.6, [
            runs(take, size=10.5, color=GRAPHITE, bold_color=INK)
        ], line_spacing=1.35)

def r_quote(prs, s):
    slide = base_slide(prs)
    accent = MS[SECTION_COLOR[s["sec"]]]
    # soft background wash + accent edge
    rect(slide, 0, 0, PW, PH, fill=BG_SOFT)
    rect(slide, 0, 0, 0.075, PH, fill=accent["base"])
    # quote mark
    text(slide, 1.15, 1.55, 1.2, 1.0, [[
        ("\u201c", {"size": 90, "bold": True, "color": accent["l50"],
         "font": "Georgia"})]])
    # the quote
    text(slide, 1.30, 2.45, PW-2.6, 2.4, [
        runs(s["quote"], size=30, color=INK, bold_color=INK)
    ], line_spacing=1.34)
    # attribution
    text(slide, 1.30, PH-2.05, PW-2.6, 0.4, [
        runs(s["cite"], size=13, color=SLATE, bold_color=INK)
    ])

def r_section_intro(prs, s):
    slide, accent, ctop = content_frame(prs, s["sec"], s["kicker"],
        runs(s["h"], size=22, color=INK, bold_color=INK),
        runs(s["desc"], size=12.5, color=SLATE),
        runs(s["src"], size=10.5, color=SLATE, bold_color=GRAPHITE),
        title_size=22)
    cy = ctop + 0.15
    ch = CONTENT_BOT - cy - 0.10
    # left: eyebrow + big statement + lead
    lw = CONTENT_W * 0.58
    text(slide, MARGIN_X, cy+0.30, lw-0.3, 0.24, [[
        (s["eyebrow"], {"size": 10, "bold": True, "color": accent["d700"],
         "tracking": 2.4, "caps": True})]])
    text(slide, MARGIN_X, cy+0.58, lw-0.3, 1.2, [
        runs(s["statement"], size=29, color=INK, bold_color=INK)
    ], line_spacing=1.16)
    text(slide, MARGIN_X, cy+1.95, lw-0.5, 1.6, [
        runs(s["lead"], size=13.5, color=SLATE)
    ], line_spacing=1.55)
    # right: aside rows with accent left border
    ax = MARGIN_X + lw + 0.3
    aw = CONTENT_W - lw - 0.3
    rect(slide, ax, cy+0.30, 0.04, ch-0.6, fill=accent["base"])
    rows = s["aside"]   # list of (k, v)
    rh = (ch - 0.6) / len(rows)
    for i, (k, v) in enumerate(rows):
        ry = cy + 0.30 + i*rh
        text(slide, ax+0.22, ry+rh*0.5-0.34, aw-0.22, 0.7, [
            [(k, {"size": 9.5, "bold": True, "color": SLATE, "tracking": 1.4, "caps": True})],
            [(v, {"size": 13, "color": INK, "space_before": 3})],
        ], line_spacing=1.35)

def r_pillar_grid(prs, s):
    slide, accent, ctop = content_frame(prs, s["sec"], s["kicker"],
        runs(s["h"], size=22, color=INK, bold_color=INK),
        runs(s["desc"], size=12.5, color=SLATE),
        runs(s["src"], size=10.5, color=SLATE, bold_color=GRAPHITE),
        title_size=22)
    pillars = s["pillars"]   # list of (no, title, desc)
    n = len(pillars)
    gap = 0.24
    cw = (CONTENT_W - gap*(n-1)) / n
    ch = 2.9
    cy = ctop + (CONTENT_BOT - ctop - ch) / 2   # vertically centered
    for i, (no, title, desc) in enumerate(pillars):
        cx = MARGIN_X + i*(cw+gap)
        rect(slide, cx, cy, cw, ch, fill=BG_CARD, line=LINE, line_w=1.0,
             rounded=True, radius=0.05)
        rect(slide, cx, cy, cw, 0.05, fill=accent["base"])
        pad = 0.28
        text(slide, cx+pad, cy+pad, cw-2*pad, 0.4, [[
            (no, {"size": 24, "bold": True, "color": accent["base"]})]])
        text(slide, cx+pad, cy+pad+0.52, cw-2*pad, 0.6, [
            runs(title, size=15.5, color=INK, bold_color=INK)
        ], line_spacing=1.18)
        text(slide, cx+pad, cy+pad+1.16, cw-2*pad, ch-pad-1.4, [
            runs(desc, size=11.5, color=SLATE)
        ], line_spacing=1.5)

def r_flow(prs, s):
    slide, accent, ctop = content_frame(prs, s["sec"], s["kicker"],
        runs(s["h"], size=22, color=INK, bold_color=INK),
        runs(s["desc"], size=12.5, color=SLATE),
        runs(s["src"], size=10.5, color=SLATE, bold_color=GRAPHITE),
        title_size=22)
    nodes = s["nodes"]   # list of (icon, title, desc)
    n = len(nodes)
    arrow_w = 0.42
    nh = 2.15
    cy = ctop + (CONTENT_BOT - ctop - nh) / 2
    nw = (CONTENT_W - arrow_w*(n-1)) / n
    node_colors = ["blue", None, "green", None, "yellow", None, "red"]
    for i, (icon, title, desc) in enumerate(nodes):
        cx = MARGIN_X + i*(nw+arrow_w)
        rect(slide, cx, cy, nw, nh, fill=BG_CARD, line=LINE, line_w=1.0,
             rounded=True, radius=0.07)
        # icon chip
        ckey = ["blue","green","yellow","red"][i % 4]
        isz = 0.46
        ix = cx + (nw-isz)/2
        rect(slide, ix, cy+0.26, isz, isz, fill=MS[ckey]["base"],
             rounded=True, radius=0.20)
        text(slide, ix, cy+0.26, isz, isz, [[
            (str(icon), {"size": 16, "bold": True,
             "color": INK if ckey=="yellow" else WHITE})]],
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        text(slide, cx+0.14, cy+0.84, nw-0.28, 0.34, [[
            (title, {"size": 13, "bold": True, "color": INK})]],
            align=PP_ALIGN.CENTER)
        text(slide, cx+0.16, cy+1.18, nw-0.32, nh-1.3, [
            runs(desc, size=10.5, color=SLATE)
        ], line_spacing=1.35, align=PP_ALIGN.CENTER)
        # arrow
        if i < n-1:
            text(slide, cx+nw, cy, arrow_w, nh, [[
                ("\u2192", {"size": 20, "color": PEARL})]],
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


# ============================================================
# RENDERERS: layer-stack, before-after, roadmap, pricing, risks-table
# ============================================================
def r_layer_stack(prs, s):
    slide, accent, ctop = content_frame(prs, s["sec"], s["kicker"],
        runs(s["h"], size=22, color=INK, bold_color=INK),
        runs(s["desc"], size=12.5, color=SLATE),
        runs(s["src"], size=10.5, color=SLATE, bold_color=GRAPHITE),
        title_size=22)
    layers = s["layers"]   # list of (tag, title, desc)
    n = len(layers)
    gap = 0.13
    cy = ctop + 0.10
    avail = CONTENT_BOT - cy - 0.05
    lh = (avail - gap*(n-1)) / n
    tag_w = 1.95
    lkeys = ["blue", "green", "yellow", "red"]
    for i, (tag, title, desc) in enumerate(layers):
        ly = cy + i*(lh+gap)
        ckey = lkeys[i % 4]
        # tag block (colored)
        rect(slide, MARGIN_X, ly, tag_w, lh, fill=MS[ckey]["base"],
             rounded=True, radius=0.05)
        text(slide, MARGIN_X+0.16, ly, tag_w-0.32, lh, [[
            (tag, {"size": 12, "bold": True,
             "color": INK if ckey == "yellow" else WHITE,
             "tracking": 0.4})]], anchor=MSO_ANCHOR.MIDDLE)
        # body block
        bx = MARGIN_X + tag_w + gap
        bw = CONTENT_W - tag_w - gap
        rect(slide, bx, ly, bw, lh, fill=BG_CARD, line=LINE, line_w=1.0,
             rounded=True, radius=0.05)
        text(slide, bx+0.26, ly+lh*0.5-0.34, bw-0.5, 0.34, [[
            (title, {"size": 14, "bold": True, "color": INK})]])
        text(slide, bx+0.26, ly+lh*0.5+0.02, bw-0.5, 0.5, [
            runs(desc, size=11, color=SLATE)
        ], line_spacing=1.3)

def r_before_after(prs, s):
    slide, accent, ctop = content_frame(prs, s["sec"], s["kicker"],
        runs(s["h"], size=22, color=INK, bold_color=INK),
        runs(s["desc"], size=12.5, color=SLATE),
        runs(s["src"], size=10.5, color=SLATE, bold_color=GRAPHITE),
        title_size=22)
    cy = ctop + 0.10
    # two columns + stat row below
    stat_h = 1.05
    col_h = CONTENT_BOT - cy - stat_h - 0.22
    arrow_w = 0.5
    colw = (CONTENT_W - arrow_w) / 2
    # before col
    def _col(x, label, items, is_after):
        edge = accent["base"] if is_after else SILVER
        rect(slide, x, cy, colw, col_h, fill=BG_CARD if is_after else BG_SOFT,
             line=LINE, line_w=1.0, rounded=True, radius=0.04)
        rect(slide, x, cy, 0.05, col_h, fill=edge)
        text(slide, x+0.26, cy+0.20, colw-0.5, 0.24, [[
            (label, {"size": 10, "bold": True,
             "color": accent["d700"] if is_after else SLATE,
             "tracking": 1.6, "caps": True})]])
        iy = cy + 0.58
        for it in items:
            text(slide, x+0.42, iy, colw-0.7, 0.5, [
                runs(it, size=11, color=GRAPHITE if is_after else SLATE)
            ], line_spacing=1.3)
            # bullet dot
            oval(slide, x+0.26, iy+0.06, 0.07, 0.07,
                 fill=edge)
            iy += 0.50
    _col(MARGIN_X, s["before_label"], s["before_items"], False)
    text(slide, MARGIN_X+colw, cy, arrow_w, col_h, [[
        ("\u2192", {"size": 22, "color": PEARL})]],
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    _col(MARGIN_X+colw+arrow_w, s["after_label"], s["after_items"], True)
    # stat row
    sy = cy + col_h + 0.18
    stats = s["stats"]   # list of (label, n, desc, featured?)
    sgap = 0.16
    sw = (CONTENT_W - sgap*(len(stats)-1)) / len(stats)
    for i, st in enumerate(stats):
        label, num, desc = st[0], st[1], st[2]
        feat = st[3] if len(st) > 3 else False
        sx = MARGIN_X + i*(sw+sgap)
        rect(slide, sx, sy, sw, stat_h, fill=accent["l50"] if feat else BG_CARD,
             line=accent["base"] if feat else LINE, line_w=1.0,
             rounded=True, radius=0.06)
        text(slide, sx+0.20, sy+0.13, sw-0.4, 0.18, [[
            (label, {"size": 8.5, "bold": True, "color": accent["d700"],
             "tracking": 1.4, "caps": True})]])
        text(slide, sx+0.20, sy+0.32, sw*0.5, 0.5, [[
            (num, {"size": 23, "bold": True, "color": INK})]])
        text(slide, sx+sw*0.46, sy+0.34, sw*0.5-0.2, 0.6, [
            runs(desc, size=9.5, color=SLATE)
        ], line_spacing=1.25, anchor=MSO_ANCHOR.MIDDLE)

def r_roadmap(prs, s):
    slide, accent, ctop = content_frame(prs, s["sec"], s["kicker"],
        runs(s["h"], size=22, color=INK, bold_color=INK),
        runs(s["desc"], size=12.5, color=SLATE),
        runs(s["src"], size=10.5, color=SLATE, bold_color=GRAPHITE),
        title_size=22)
    phases = s["phases"]   # (no, name, when, items[], exit_label, exit_text)
    n = len(phases)
    gap = 0.18
    cw = (CONTENT_W - gap*(n-1)) / n
    cy = ctop + 0.10
    ch = CONTENT_BOT - cy - 0.05
    pkeys = ["blue", "green", "yellow", "red"]
    for i, (no, name, when, items, ex_lbl, ex_text) in enumerate(phases):
        cx = MARGIN_X + i*(cw+gap)
        ckey = pkeys[i % 4]
        rect(slide, cx, cy, cw, ch, fill=BG_CARD, line=LINE, line_w=1.0,
             rounded=True, radius=0.045)
        # colored header band
        hdr_h = 0.62
        rect(slide, cx, cy, cw, hdr_h, fill=MS[ckey]["base"], rounded=True, radius=0.10)
        rect(slide, cx, cy+hdr_h*0.5, cw, hdr_h*0.5, fill=MS[ckey]["base"])  # square bottom
        on_y = ckey == "yellow"
        text(slide, cx+0.20, cy+0.09, cw-0.4, 0.20, [[
            (no, {"size": 8.5, "bold": True,
             "color": INK if on_y else WHITE, "tracking": 1.2, "caps": True})]])
        text(slide, cx+0.20, cy+0.27, cw-0.4, 0.30, [[
            (name, {"size": 15, "bold": True, "color": INK if on_y else WHITE})]])
        # when
        text(slide, cx+0.20, cy+hdr_h+0.12, cw-0.4, 0.22, [[
            (when, {"size": 9.5, "bold": True, "color": SLATE,
             "tracking": 0.8, "caps": True})]])
        # items
        iy = cy + hdr_h + 0.40
        for it in items:
            oval(slide, cx+0.22, iy+0.05, 0.06, 0.06, fill=MS[ckey]["base"])
            text(slide, cx+0.38, iy, cw-0.58, 0.5, [
                runs(it, size=10, color=GRAPHITE)
            ], line_spacing=1.25)
            iy += 0.42
        # exit criteria pinned to bottom
        ex_h = 1.15
        ex_y = cy + ch - ex_h
        hline(slide, cx+0.20, ex_y, cw-0.4, LINE2, 0.75)
        text(slide, cx+0.20, ex_y+0.10, cw-0.4, 0.18, [[
            (ex_lbl, {"size": 8, "bold": True, "color": accent["d700"],
             "tracking": 1.2, "caps": True})]])
        text(slide, cx+0.20, ex_y+0.30, cw-0.4, ex_h-0.4, [
            runs(ex_text, size=10, color=SLATE, bold_color=INK)
        ], line_spacing=1.32)

def r_pricing(prs, s):
    slide, accent, ctop = content_frame(prs, s["sec"], s["kicker"],
        runs(s["h"], size=22, color=INK, bold_color=INK),
        runs(s["desc"], size=12.5, color=SLATE),
        runs(s["src"], size=10.5, color=SLATE, bold_color=GRAPHITE),
        title_size=22)
    tiers = s["tiers"]   # (name, price, unit, sub, items[], foot, featured, badge)
    n = len(tiers)
    gap = 0.20
    cw = (CONTENT_W - gap*(n-1)) / n
    cy = ctop + 0.10
    ch = CONTENT_BOT - cy - 0.05
    tkeys = ["blue", "green", "red"]
    for i, (name, price, unit, sub, items, foot, featured, badge) in enumerate(tiers):
        cx = MARGIN_X + i*(cw+gap)
        ckey = tkeys[i % 3]
        rect(slide, cx, cy, cw, ch, fill=BG_CARD,
             line=accent["base"] if featured else LINE,
             line_w=2.0 if featured else 1.0, rounded=True, radius=0.05)
        pad = 0.24
        yy = cy + pad
        if badge:
            bw = 0.30 + len(badge)*0.072
            rect(slide, cx+pad, yy, bw, 0.24, fill=accent["base"], rounded=True, radius=0.5)
            text(slide, cx+pad, yy, bw, 0.24, [[
                (badge, {"size": 7.5, "bold": True, "color": WHITE,
                 "tracking": 0.6, "caps": True})]],
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            yy += 0.34
        text(slide, cx+pad, yy, cw-2*pad, 0.20, [[
            (name, {"size": 9.5, "bold": True, "color": MS[ckey]["d700"],
             "tracking": 1.4, "caps": True})]])
        text(slide, cx+pad, yy+0.22, cw-2*pad, 0.5, [[
            (price, {"size": 30, "bold": True, "color": INK}),
            (" " + unit, {"size": 12, "color": SLATE}),
        ]])
        text(slide, cx+pad, yy+0.72, cw-2*pad, 0.24, [[
            (sub, {"size": 10.5, "color": SLATE})]])
        hline(slide, cx+pad, yy+1.02, cw-2*pad, LINE2, 0.75)
        iy = yy + 1.16
        for it in items:
            text(slide, cx+pad+0.18, iy, cw-2*pad-0.18, 0.4, [
                runs(it, size=10.5, color=GRAPHITE)
            ], line_spacing=1.25)
            text(slide, cx+pad, iy, 0.16, 0.24, [[
                ("+", {"size": 12, "bold": True, "color": accent["base"]})]])
            iy += 0.34
        # foot pinned to bottom
        foot_h = 0.78
        fy = cy + ch - foot_h
        rect(slide, cx, fy, cw, foot_h, fill=BG_SOFT)
        rect(slide, cx, fy, cw, foot_h, fill=None, line=LINE2, line_w=0.75)
        text(slide, cx+pad, fy+0.12, cw-2*pad, foot_h-0.24, [
            runs(foot, size=10, color=SLATE, bold_color=INK)
        ], line_spacing=1.32)

def r_risks_table(prs, s):
    """Risk table: Risk | Likelihood | Impact | Mitigation, with tag pills."""
    slide, accent, ctop = content_frame(prs, s["sec"], s["kicker"],
        runs(s["h"], size=22, color=INK, bold_color=INK),
        runs(s["desc"], size=12.5, color=SLATE),
        runs(s["src"], size=10.5, color=SLATE, bold_color=GRAPHITE),
        title_size=22)
    rows = s["rows"]   # list of (risk, likelihood, impact, mitigation)
    col_w = [0.24*CONTENT_W, 0.14*CONTENT_W, 0.14*CONTENT_W, 0.48*CONTENT_W]
    col_x = [MARGIN_X]
    for cw in col_w[:-1]:
        col_x.append(col_x[-1]+cw)
    HEADER_H = 0.40
    LINE_H = 0.205
    PAD_V = 0.20
    body_h = []
    for risk, lk, im, mit in rows:
        lines = max(_est_lines(risk, col_w[0]-0.22, 11.5),
                    _est_lines(mit, col_w[3]-0.22, 11))
        body_h.append(max(0.46, lines*LINE_H + PAD_V))
    ty = ctop
    nrows = len(rows)+1
    tbl = native_table(slide, MARGIN_X, ty, CONTENT_W, nrows, 4,
                       col_widths=col_w, row_h=0.40)
    tbl.rows[0].height = IN(HEADER_H)
    for ri, rh in enumerate(body_h, start=1):
        tbl.rows[ri].height = IN(rh)
    for ci, label in enumerate(["Risk", "Likelihood", "Impact", "Mitigation"]):
        cell_text(tbl.cell(0, ci), [(label, {})], size=9, bold=True, color=SLATE,
                  fill=WHITE, anchor=MSO_ANCHOR.BOTTOM)
    for ri, (risk, lk, im, mit) in enumerate(rows, start=1):
        zebra = BG_SOFT if ri % 2 == 0 else WHITE
        cell_text(tbl.cell(ri,0), [(risk, {"bold": True, "color": INK})], size=11.5, fill=zebra)
        cell_text(tbl.cell(ri,1), "", fill=zebra)
        cell_text(tbl.cell(ri,2), "", fill=zebra)
        cell_text(tbl.cell(ri,3), [(mit, {})], size=11, color=GRAPHITE, fill=zebra)
    hline(slide, MARGIN_X, ty+HEADER_H, CONTENT_W, GRAPHITE, 1.5)
    # overlay tag pills
    y_cursor = ty + HEADER_H
    for ri, (risk, lk, im, mit) in enumerate(rows):
        rh = body_h[ri]
        for ci, tag in [(1, lk), (2, im)]:
            px = col_x[ci] + 0.10
            py = y_cursor + (rh-0.26)/2
            _tag_pill(slide, px, py, tag)
        y_cursor += rh
    table_bottom = ty + HEADER_H + sum(body_h)
    # pull quote
    if s.get("pull"):
        py = table_bottom + 0.20
        ph = CONTENT_BOT - py
        if ph > 0.70:
            lbl, txt, cite = s["pull"]
            rect(slide, MARGIN_X, py, CONTENT_W, ph, fill=accent["l50"],
                 rounded=True, radius=0.04)
            rect(slide, MARGIN_X, py, 0.045, ph, fill=accent["base"])
            text(slide, MARGIN_X+0.24, py+0.13, CONTENT_W-0.5, 0.18, [[
                (lbl, {"size": 9, "bold": True, "color": accent["d700"],
                 "tracking": 1.8, "caps": True})]])
            text(slide, MARGIN_X+0.24, py+0.34, CONTENT_W-0.6, ph-0.62, [
                runs(txt, size=12.5, color=INK, bold_color=INK)
            ], line_spacing=1.32)
            text(slide, MARGIN_X+0.24, py+ph-0.26, CONTENT_W-0.5, 0.20, [
                runs(cite, size=9.5, color=SLATE, bold_color=GRAPHITE)
            ])


# ============================================================
# RENDERERS: quadrant, agenda, recommendation, team, faq, dashboard
# ============================================================
def _vlabel(slide, x, y, h, txt):
    """A vertical axis label, rotated 270 degrees, centered on its band."""
    box = slide.shapes.add_textbox(IN(x), IN(y), IN(h), IN(0.26))
    box.rotation = 270
    # after rotation, reposition so it sits along the left edge
    box.left = IN(x - h/2 + 0.13)
    box.top = IN(y + h/2 - 0.13)
    tf = box.text_frame
    tf.word_wrap = False
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = txt
    f = r.font
    f.size = Pt(8.5); f.bold = True; f.name = SANS
    f.color.rgb = SLATE
    r._r.get_or_add_rPr().set("spc", "100")
    return box

def r_quadrant(prs, s):
    slide, accent, ctop = content_frame(prs, s["sec"], s["kicker"],
        runs(s["h"], size=22, color=INK, bold_color=INK),
        runs(s["desc"], size=12.5, color=SLATE),
        runs(s["src"], size=10.5, color=SLATE, bold_color=GRAPHITE),
        title_size=22)
    cy = ctop + 0.10
    legend_h = 0.34
    qh = CONTENT_BOT - cy - legend_h - 0.30
    yaxis_w = 0.34
    qx = MARGIN_X + yaxis_w
    qw = CONTENT_W - yaxis_w
    xaxis_h = 0.26
    grid_h = qh - xaxis_h
    # y-axis label, properly rotated
    _vlabel(slide, MARGIN_X, cy, grid_h, s["yaxis"])
    # 2x2 grid
    cells = s["cells"]
    half_w = qw/2
    half_h = grid_h/2
    for idx, (lbl, title, p, hot) in enumerate(cells):
        r_, c_ = idx//2, idx%2
        cellx = qx + c_*half_w
        celly = cy + r_*half_h
        rect(slide, cellx+0.01, celly+0.01, half_w-0.02, half_h-0.02,
             fill=accent["l50"] if hot else BG_CARD, line=LINE, line_w=1.0)
        text(slide, cellx+0.18, celly+0.14, half_w-0.36, 0.18, [[
            (lbl, {"size": 8, "bold": True,
             "color": accent["d700"] if hot else SILVER,
             "tracking": 1.0, "caps": True})]])
        text(slide, cellx+0.18, celly+0.34, half_w-0.36, 0.3, [[
            (title, {"size": 13, "bold": True, "color": INK})]])
        text(slide, cellx+0.18, celly+0.64, half_w-0.36, half_h-0.78, [
            runs(p, size=10.5, color=SLATE)
        ], line_spacing=1.35)
    # dots
    for kind, lft, top in s["dots"]:
        dot_c = accent["base"] if kind == "us" else SILVER
        dx = qx + float(lft.rstrip("%"))/100*qw - 0.08
        dy = cy + float(top.rstrip("%"))/100*grid_h - 0.08
        oval(slide, dx, dy, 0.16, 0.16, fill=dot_c, line=WHITE, line_w=2.0)
    # x-axis label
    text(slide, qx, cy+grid_h+0.04, qw, xaxis_h, [[
        (s["xaxis"], {"size": 8.5, "bold": True, "color": SLATE,
         "tracking": 1.0, "caps": True})]], align=PP_ALIGN.CENTER)
    # legend
    lx = MARGIN_X
    ly = cy + qh + 0.12
    for color_kind, label in s["legend"]:
        col = MS["red"]["base"] if "red" in color_kind else SILVER
        oval(slide, lx, ly+0.02, 0.13, 0.13, fill=col)
        tw = 0.22 + len(label)*0.085
        text(slide, lx+0.20, ly-0.04, tw, 0.24, [[
            (label, {"size": 10.5, "color": SLATE})]], anchor=MSO_ANCHOR.MIDDLE)
        lx += 0.20 + tw + 0.25

def r_agenda(prs, s):
    slide, accent, ctop = content_frame(prs, s["sec"], s["kicker"],
        runs(s["h"], size=22, color=INK, bold_color=INK),
        runs(s["desc"], size=12.5, color=SLATE),
        runs(s["src"], size=10.5, color=SLATE, bold_color=GRAPHITE),
        title_size=22)
    secs = s["sections"]   # list of (no, name, range, items[])
    n = len(secs)
    cols = 3
    rows_ct = (n + cols - 1)//cols
    gap = 0.20
    cw = (CONTENT_W - gap*(cols-1))/cols
    cy = ctop + 0.12
    avail = CONTENT_BOT - cy - 0.05
    chh = (avail - gap*(rows_ct-1))/rows_ct
    skeys = ["blue","green","yellow","red","blue","green"]
    for i, (no, name, rng, items) in enumerate(secs):
        r_, c_ = i//cols, i%cols
        cx = MARGIN_X + c_*(cw+gap)
        cyy = cy + r_*(chh+gap)
        ckey = skeys[i % 6]
        rect(slide, cx, cyy, cw, chh, fill=BG_CARD, line=LINE, line_w=1.0,
             rounded=True, radius=0.05)
        rect(slide, cx, cyy, cw, 0.05, fill=MS[ckey]["base"])
        pad = 0.22
        text(slide, cx+pad, cyy+pad, cw-2*pad, 0.20, [[
            (no, {"size": 11, "bold": True, "color": SILVER, "tracking": 1.0})]])
        text(slide, cx+pad, cyy+pad+0.20, cw-2*pad, 0.32, [[
            (name, {"size": 14, "bold": True, "color": INK})]])
        text(slide, cx+pad, cyy+pad+0.52, cw-2*pad, 0.18, [[
            (rng, {"size": 8.5, "bold": True, "color": SLATE, "tracking": 1.0, "caps": True})]])
        iy = cyy + pad + 0.78
        for it in items:
            oval(slide, cx+pad, iy+0.05, 0.05, 0.05, fill=PEARL)
            text(slide, cx+pad+0.14, iy, cw-2*pad-0.14, 0.26, [[
                (it, {"size": 10, "color": SLATE})]])
            iy += 0.24

def r_recommendation(prs, s):
    slide, accent, ctop = content_frame(prs, s["sec"], s["kicker"],
        runs(s["h"], size=22, color=INK, bold_color=INK),
        runs(s["desc"], size=12.5, color=SLATE),
        runs(s["src"], size=10.5, color=SLATE, bold_color=GRAPHITE),
        title_size=22)
    cy = ctop + 0.12
    ch = CONTENT_BOT - cy - 0.05
    # left: 3 stacked pillars
    lw = CONTENT_W * 0.58
    pillars = s["pillars"]   # (no, title, desc)
    pgap = 0.14
    ph = (ch - pgap*(len(pillars)-1))/len(pillars)
    for i, (no, title, desc) in enumerate(pillars):
        py = cy + i*(ph+pgap)
        rect(slide, MARGIN_X, py, lw, ph, fill=BG_CARD, line=LINE, line_w=1.0,
             rounded=True, radius=0.06)
        rect(slide, MARGIN_X, py, 0.05, ph, fill=accent["base"])
        text(slide, MARGIN_X+0.26, py+0.16, 0.6, 0.3, [[
            (no, {"size": 17, "bold": True, "color": accent["base"]})]])
        text(slide, MARGIN_X+0.74, py+0.16, lw-1.0, 0.3, [[
            (title, {"size": 13.5, "bold": True, "color": INK})]])
        text(slide, MARGIN_X+0.74, py+0.46, lw-1.0, ph-0.6, [
            runs(desc, size=10.5, color=SLATE)
        ], line_spacing=1.35)
    # right: dark ask card + next-step pull
    rx = MARGIN_X + lw + 0.28
    rw = CONTENT_W - lw - 0.28
    ask_h = ch * 0.52
    rect(slide, rx, cy, rw, ask_h, fill=GRAPHITE, rounded=True, radius=0.05)
    text(slide, rx+0.26, cy+0.22, rw-0.52, 0.2, [[
        (s["ask_label"], {"size": 9, "bold": True, "color": RGBColor(0xC0,0xC0,0xC0),
         "tracking": 1.8, "caps": True})]])
    text(slide, rx+0.26, cy+0.46, rw-0.52, 0.7, [[
        (s["ask_n"], {"size": 40, "bold": True, "color": WHITE})]])
    text(slide, rx+0.26, cy+1.20, rw-0.52, ask_h-1.4, [
        runs(s["ask_desc"], size=11, color=RGBColor(0xDD,0xDD,0xDD), bold_color=WHITE)
    ], line_spacing=1.4)
    # next-step pull
    py = cy + ask_h + 0.16
    ph2 = ch - ask_h - 0.16
    rect(slide, rx, py, rw, ph2, fill=accent["l50"], rounded=True, radius=0.05)
    rect(slide, rx, py, 0.045, ph2, fill=accent["base"])
    text(slide, rx+0.22, py+0.16, rw-0.4, 0.18, [[
        (s["next_label"], {"size": 9, "bold": True, "color": accent["d700"],
         "tracking": 1.8, "caps": True})]])
    text(slide, rx+0.22, py+0.38, rw-0.42, ph2-0.8, [
        runs(s["next_text"], size=11.5, color=INK, bold_color=INK)
    ], line_spacing=1.4)
    text(slide, rx+0.22, py+ph2-0.28, rw-0.4, 0.20, [
        runs(s["contact"], size=9.5, color=SLATE, bold_color=GRAPHITE)
    ])

def r_team(prs, s):
    slide, accent, ctop = content_frame(prs, s["sec"], s["kicker"],
        runs(s["h"], size=22, color=INK, bold_color=INK),
        runs(s["desc"], size=12.5, color=SLATE),
        runs(s["src"], size=10.5, color=SLATE, bold_color=GRAPHITE),
        title_size=22)
    people = s["people"]   # (initials, name, role, desc)
    n = len(people)
    gap = 0.22
    cw = (CONTENT_W - gap*(n-1))/n
    chh = 2.85
    cy = ctop + (CONTENT_BOT - ctop - chh)/2
    pkeys = ["blue","green","yellow","red"]
    for i, (initials, name, role, desc) in enumerate(people):
        cx = MARGIN_X + i*(cw+gap)
        ckey = pkeys[i % 4]
        rect(slide, cx, cy, cw, chh, fill=BG_CARD, line=LINE, line_w=1.0,
             rounded=True, radius=0.06)
        pad = 0.26
        # avatar circle
        av = 0.62
        oval(slide, cx+pad, cy+pad, av, av, fill=MS[ckey]["base"])
        text(slide, cx+pad, cy+pad, av, av, [[
            (initials, {"size": 18, "bold": True,
             "color": INK if ckey=="yellow" else WHITE})]],
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        text(slide, cx+pad, cy+pad+0.78, cw-2*pad, 0.3, [[
            (name, {"size": 13.5, "bold": True, "color": INK})]])
        text(slide, cx+pad, cy+pad+1.08, cw-2*pad, 0.2, [[
            (role, {"size": 9, "bold": True, "color": accent["d700"],
             "tracking": 1.0, "caps": True})]])
        text(slide, cx+pad, cy+pad+1.34, cw-2*pad, chh-pad-1.5, [
            runs(desc, size=10.5, color=SLATE)
        ], line_spacing=1.4)

def r_faq(prs, s):
    slide, accent, ctop = content_frame(prs, s["sec"], s["kicker"],
        runs(s["h"], size=22, color=INK, bold_color=INK),
        runs(s["desc"], size=12.5, color=SLATE),
        runs(s["src"], size=10.5, color=SLATE, bold_color=GRAPHITE),
        title_size=22)
    items = s["items"]   # (q, a)
    n = len(items)
    cols = 2
    rows_ct = (n+cols-1)//cols
    cgap = 0.45
    rgap = 0.28
    cw = (CONTENT_W - cgap)/cols
    cy = ctop + 0.15
    avail = CONTENT_BOT - cy - 0.05
    rh = (avail - rgap*(rows_ct-1))/rows_ct
    for i, (q, a) in enumerate(items):
        r_, c_ = i//cols, i%cols
        cx = MARGIN_X + c_*(cw+cgap)
        cyy = cy + r_*(rh+rgap)
        # Q marker
        text(slide, cx, cyy, 0.26, 0.3, [[
            ("Q", {"size": 13, "bold": True, "color": accent["base"]})]])
        text(slide, cx+0.30, cyy, cw-0.30, 0.5, [
            runs(q, size=12.5, color=INK, bold_color=INK)
        ], line_spacing=1.2)
        text(slide, cx+0.30, cyy+0.42, cw-0.30, rh-0.5, [
            runs(a, size=11, color=SLATE, bold_color=GRAPHITE)
        ], line_spacing=1.42)

def r_metrics_dashboard(prs, s):
    slide, accent, ctop = content_frame(prs, s["sec"], s["kicker"],
        runs(s["h"], size=22, color=INK, bold_color=INK),
        runs(s["desc"], size=12.5, color=SLATE),
        runs(s["src"], size=10.5, color=SLATE, bold_color=GRAPHITE),
        title_size=22)
    cy = ctop + 0.12
    ch = CONTENT_BOT - cy - 0.05
    # featured tile (left, tall) + 4 small tiles (2x2 right)
    feat = s["featured"]   # (label, n, unit, desc)
    tiles = s["tiles"]     # list of (label, n, unit, desc)
    fw = CONTENT_W * 0.30
    rect(slide, MARGIN_X, cy, fw, ch, fill=GRAPHITE, rounded=True, radius=0.05)
    text(slide, MARGIN_X+0.26, cy+0.30, fw-0.52, 0.2, [[
        (feat[0], {"size": 9, "bold": True, "color": RGBColor(0xB8,0xB8,0xB8),
         "tracking": 1.6, "caps": True})]])
    text(slide, MARGIN_X+0.26, cy+0.58, fw-0.52, 1.0, [[
        (feat[1], {"size": 56, "bold": True, "color": WHITE}),
        (feat[2], {"size": 20, "color": accent["base"]}),
    ]])
    text(slide, MARGIN_X+0.26, cy+ch-1.2, fw-0.52, 1.0, [
        runs(feat[3], size=11.5, color=RGBColor(0xDD,0xDD,0xDD), bold_color=WHITE)
    ], line_spacing=1.42)
    # right grid 2x2
    gx = MARGIN_X + fw + 0.22
    gw = CONTENT_W - fw - 0.22
    gap = 0.18
    tw = (gw - gap)/2
    th = (ch - gap)/2
    for i, (label, num, unit, desc) in enumerate(tiles[:4]):
        r_, c_ = i//2, i%2
        tx = gx + c_*(tw+gap)
        tyy = cy + r_*(th+gap)
        rect(slide, tx, tyy, tw, th, fill=BG_CARD, line=LINE, line_w=1.0,
             rounded=True, radius=0.06)
        text(slide, tx+0.22, tyy+0.20, tw-0.44, 0.18, [[
            (label, {"size": 8.5, "bold": True, "color": accent["d700"],
             "tracking": 1.2, "caps": True})]])
        text(slide, tx+0.22, tyy+0.42, tw-0.44, 0.6, [[
            (num, {"size": 34, "bold": True, "color": INK}),
            (unit, {"size": 15, "color": SLATE}),
        ]])
        text(slide, tx+0.22, tyy+th-0.5, tw-0.44, 0.44, [
            runs(desc, size=10, color=SLATE)
        ], line_spacing=1.3)


# ============================================================
# RENDERERS: case-study, architecture-detail, vendor-compare, raci,
#             gantt, decision-matrix, stakeholder-map, org-chart,
#             appendix-divider, contact
# ============================================================
def r_case_study(prs, s):
    slide, accent, ctop = content_frame(prs, s["sec"], s["kicker"],
        runs(s["h"], size=22, color=INK, bold_color=INK),
        runs(s["desc"], size=12.5, color=SLATE),
        runs(s["src"], size=10.5, color=SLATE, bold_color=GRAPHITE),
        title_size=22)
    cy = ctop + 0.12
    ch = CONTENT_BOT - cy - 0.05
    hw = CONTENT_W * 0.32
    rect(slide, MARGIN_X, cy, hw, ch, fill=GRAPHITE, rounded=True, radius=0.05)
    pad = 0.26
    text(slide, MARGIN_X+pad, cy+pad, hw-2*pad, 0.18, [[
        (s["hero_tag"], {"size": 9, "bold": True, "color": accent["base"],
         "tracking": 1.6, "caps": True})]])
    text(slide, MARGIN_X+pad, cy+pad+0.24, hw-2*pad, 0.8, [
        runs(s["hero_name"], size=20, color=WHITE, bold_color=WHITE)
    ], line_spacing=1.15)
    text(slide, MARGIN_X+pad, cy+pad+1.05, hw-2*pad, 0.9, [
        runs(s["hero_meta"], size=10.5, color=RGBColor(0xB0,0xB0,0xB0))
    ], line_spacing=1.45)
    text(slide, MARGIN_X+pad, cy+ch-1.18, hw-2*pad, 0.6, [[
        (s["hero_stat_n"], {"size": 36, "bold": True, "color": WHITE})]])
    text(slide, MARGIN_X+pad, cy+ch-0.62, hw-2*pad, 0.5, [
        runs(s["hero_stat_l"], size=10, color=RGBColor(0xC0,0xC0,0xC0))
    ], line_spacing=1.3)
    sx = MARGIN_X + hw + 0.26
    sw = CONTENT_W - hw - 0.26
    steps = s["steps"]
    sgap = 0.14
    sh = (ch - sgap*(len(steps)-1))/len(steps)
    skeys = ["red", "yellow", "green"]
    for i, (kind, txt) in enumerate(steps):
        syy = cy + i*(sh+sgap)
        rect(slide, sx, syy, sw, sh, fill=BG_CARD, line=LINE, line_w=1.0,
             rounded=True, radius=0.06)
        text(slide, sx+0.24, syy, 1.1, sh, [[
            (kind, {"size": 10, "bold": True, "color": MS[skeys[i%3]]["d700"],
             "tracking": 0.8, "caps": True})]], anchor=MSO_ANCHOR.MIDDLE)
        text(slide, sx+1.35, syy+0.14, sw-1.6, sh-0.28, [
            runs(txt, size=11.5, color=GRAPHITE, bold_color=INK)
        ], line_spacing=1.4, anchor=MSO_ANCHOR.MIDDLE)

def r_architecture_detail(prs, s):
    slide, accent, ctop = content_frame(prs, s["sec"], s["kicker"],
        runs(s["h"], size=22, color=INK, bold_color=INK),
        runs(s["desc"], size=12.5, color=SLATE),
        runs(s["src"], size=10.5, color=SLATE, bold_color=GRAPHITE),
        title_size=22)
    cy = ctop + 0.12
    ch = CONTENT_BOT - cy - 0.05
    cw_l = CONTENT_W * 0.66
    rect(slide, MARGIN_X, cy, cw_l, ch, fill=BG_SOFT, line=LINE, line_w=1.0,
         rounded=True, radius=0.04)
    tiers = s["tiers"]
    pad = 0.22
    tgap = 0.16
    th = (ch - 2*pad - tgap*(len(tiers)-1))/len(tiers)
    for ti, (tlabel, nodes) in enumerate(tiers):
        tyy = cy + pad + ti*(th+tgap)
        text(slide, MARGIN_X+pad, tyy, cw_l-2*pad, 0.16, [[
            (tlabel, {"size": 8, "bold": True, "color": SLATE,
             "tracking": 1.2, "caps": True})]])
        nrow_y = tyy + 0.20
        nrow_h = th - 0.20
        ngap = 0.10
        nw = (cw_l - 2*pad - ngap*(len(nodes)-1))/len(nodes)
        for ni, (ntitle, nsub, ckey) in enumerate(nodes):
            nx = MARGIN_X + pad + ni*(nw+ngap)
            rect(slide, nx, nrow_y, nw, nrow_h, fill=WHITE, line=LINE, line_w=1.0,
                 rounded=True, radius=0.08)
            rect(slide, nx, nrow_y, nw, 0.04, fill=MS[ckey]["base"])
            text(slide, nx+0.12, nrow_y+0.10, nw-0.24, 0.24, [[
                (ntitle, {"size": 10.5, "bold": True, "color": INK})]])
            text(slide, nx+0.12, nrow_y+0.32, nw-0.24, nrow_h-0.4, [
                runs(nsub, size=9, color=SLATE)
            ], line_spacing=1.25)
    nx = MARGIN_X + cw_l + 0.26
    nw_ = CONTENT_W - cw_l - 0.26
    notes = s["notes"]
    ngap = 0.20
    nh = (ch - ngap*(len(notes)-1))/len(notes)
    for i, (k, v) in enumerate(notes):
        nyy = cy + i*(nh+ngap)
        rect(slide, nx, nyy+0.04, 0.10, 0.10, fill=accent["base"])
        text(slide, nx+0.22, nyy, nw_-0.22, 0.24, [[
            (k, {"size": 11, "bold": True, "color": INK})]])
        text(slide, nx+0.22, nyy+0.24, nw_-0.22, nh-0.3, [
            runs(v, size=10.5, color=SLATE)
        ], line_spacing=1.4)

def r_vendor_compare(prs, s):
    slide, accent, ctop = content_frame(prs, s["sec"], s["kicker"],
        runs(s["h"], size=22, color=INK, bold_color=INK),
        runs(s["desc"], size=12.5, color=SLATE),
        runs(s["src"], size=10.5, color=SLATE, bold_color=GRAPHITE),
        title_size=22)
    cap_label = s["cap_label"]
    options = s["options"]
    rows = s["rows"]
    ncols = 1 + len(options)
    cap_w = 0.30 * CONTENT_W
    opt_w = (CONTENT_W - cap_w)/len(options)
    col_w = [cap_w] + [opt_w]*len(options)
    ty = ctop
    nrows = len(rows)+1
    HEADER_H = 0.46
    body_rh = min(0.46, (CONTENT_BOT - ty - HEADER_H - 0.7)/len(rows))
    tbl = native_table(slide, MARGIN_X, ty, CONTENT_W, nrows, ncols,
                       col_widths=col_w, row_h=0.40)
    tbl.rows[0].height = IN(HEADER_H)
    for ri in range(1, nrows):
        tbl.rows[ri].height = IN(body_rh)
    cell_text(tbl.cell(0,0), [(cap_label, {})], size=9, bold=True, color=SLATE,
              fill=WHITE, anchor=MSO_ANCHOR.BOTTOM)
    for ci, (name, win) in enumerate(options, start=1):
        cell_text(tbl.cell(0,ci), [(name, {})], size=10, bold=True,
                  color=accent["d700"] if win else INK,
                  align=PP_ALIGN.CENTER, fill=WHITE, anchor=MSO_ANCHOR.BOTTOM)
    sym = {"yes": ("\u2713", MS["green"]["d700"]),
           "no": ("\u2014", PEARL),
           "partial": ("~", MS["yellow"]["d700"])}
    for ri, (cap, vals) in enumerate(rows, start=1):
        zebra = BG_SOFT if ri % 2 == 0 else WHITE
        cell_text(tbl.cell(ri,0), [(cap, {"bold": True, "color": INK})],
                  size=11, fill=zebra)
        for ci, v in enumerate(vals, start=1):
            win = options[ci-1][1]
            cellfill = accent["l50"] if win else zebra
            if v in sym:
                glyph, col = sym[v]
                cell_text(tbl.cell(ri,ci), [(glyph, {"bold": True, "color": col})],
                          size=13, align=PP_ALIGN.CENTER, fill=cellfill)
            else:
                cell_text(tbl.cell(ri,ci), [(v, {})], size=10, color=GRAPHITE,
                          align=PP_ALIGN.CENTER, fill=cellfill)
    hline(slide, MARGIN_X, ty+HEADER_H, CONTENT_W, GRAPHITE, 1.5)
    if s.get("pull"):
        table_bottom = ty + HEADER_H + body_rh*len(rows)
        py = table_bottom + 0.18
        ph = CONTENT_BOT - py
        if ph > 0.6:
            lbl, txt, cite = s["pull"]
            rect(slide, MARGIN_X, py, CONTENT_W, ph, fill=accent["l50"], rounded=True, radius=0.04)
            rect(slide, MARGIN_X, py, 0.045, ph, fill=accent["base"])
            text(slide, MARGIN_X+0.24, py+0.12, CONTENT_W-0.5, 0.18, [[
                (lbl, {"size": 9, "bold": True, "color": accent["d700"],
                 "tracking": 1.8, "caps": True})]])
            text(slide, MARGIN_X+0.24, py+0.33, CONTENT_W-0.6, ph-0.58, [
                runs(txt, size=12, color=INK, bold_color=INK)
            ], line_spacing=1.3)

def r_raci(prs, s):
    slide, accent, ctop = content_frame(prs, s["sec"], s["kicker"],
        runs(s["h"], size=22, color=INK, bold_color=INK),
        runs(s["desc"], size=12.5, color=SLATE),
        runs(s["src"], size=10.5, color=SLATE, bold_color=GRAPHITE),
        title_size=22)
    roles = s["roles"]
    rows = s["rows"]
    ncols = 1 + len(roles)
    ws_w = 0.32 * CONTENT_W
    role_w = (CONTENT_W - ws_w)/len(roles)
    col_w = [ws_w] + [role_w]*len(roles)
    col_x = [MARGIN_X]
    for cw in col_w[:-1]:
        col_x.append(col_x[-1]+cw)
    ty = ctop
    legend_h = 0.40
    nrows = len(rows)+1
    HEADER_H = 0.42
    body_rh = min(0.50, (CONTENT_BOT - ty - HEADER_H - legend_h - 0.2)/len(rows))
    tbl = native_table(slide, MARGIN_X, ty, CONTENT_W, nrows, ncols,
                       col_widths=col_w, row_h=0.40)
    tbl.rows[0].height = IN(HEADER_H)
    for ri in range(1, nrows):
        tbl.rows[ri].height = IN(body_rh)
    cell_text(tbl.cell(0,0), [("Workstream", {})], size=9, bold=True, color=SLATE,
              fill=WHITE, anchor=MSO_ANCHOR.BOTTOM)
    for ci, role in enumerate(roles, start=1):
        cell_text(tbl.cell(0,ci), [(role, {})], size=8.5, bold=True, color=SLATE,
                  align=PP_ALIGN.CENTER, fill=WHITE, anchor=MSO_ANCHOR.BOTTOM)
    for ri, (ws, codes) in enumerate(rows, start=1):
        zebra = BG_SOFT if ri % 2 == 0 else WHITE
        cell_text(tbl.cell(ri,0), [(ws, {"bold": True, "color": INK})], size=11, fill=zebra)
        for ci in range(1, ncols):
            cell_text(tbl.cell(ri,ci), "", fill=zebra)
    hline(slide, MARGIN_X, ty+HEADER_H, CONTENT_W, GRAPHITE, 1.5)
    badge_color = {"R": "red", "A": "blue", "C": "green", "I": "yellow"}
    y_cursor = ty + HEADER_H
    for ri, (ws, codes) in enumerate(rows):
        for ci, code in enumerate(codes, start=1):
            bx = col_x[ci] + role_w/2 - 0.13
            by = y_cursor + body_rh/2 - 0.13
            ck = badge_color.get(code, "blue")
            rect(slide, bx, by, 0.26, 0.26, fill=MS[ck]["base"], rounded=True, radius=0.22)
            text(slide, bx, by, 0.26, 0.26, [[
                (code, {"size": 11, "bold": True,
                 "color": INK if ck == "yellow" else WHITE})]],
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        y_cursor += body_rh
    ly = ty + HEADER_H + body_rh*len(rows) + 0.14
    lx = MARGIN_X
    for code, label in [("R","Responsible"),("A","Accountable"),
                        ("C","Consulted"),("I","Informed")]:
        ck = badge_color[code]
        rect(slide, lx, ly, 0.24, 0.24, fill=MS[ck]["base"], rounded=True, radius=0.22)
        text(slide, lx, ly, 0.24, 0.24, [[
            (code, {"size": 10, "bold": True,
             "color": INK if ck=="yellow" else WHITE})]],
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        tw = 0.18 + len(label)*0.075
        text(slide, lx+0.32, ly, tw, 0.24, [[
            (label, {"size": 10, "color": SLATE})]], anchor=MSO_ANCHOR.MIDDLE)
        lx += 0.32 + tw + 0.30

def r_gantt(prs, s):
    slide, accent, ctop = content_frame(prs, s["sec"], s["kicker"],
        runs(s["h"], size=22, color=INK, bold_color=INK),
        runs(s["desc"], size=12.5, color=SLATE),
        runs(s["src"], size=10.5, color=SLATE, bold_color=GRAPHITE),
        title_size=22)
    cy = ctop + 0.20
    ch = CONTENT_BOT - cy - 0.10
    label_w = 1.95
    track_x = MARGIN_X + label_w
    track_w = CONTENT_W - label_w
    units = s["units"]
    bars = s["bars"]
    text(slide, MARGIN_X, cy, label_w, 0.24, [[
        ("Workstream", {"size": 8, "bold": True, "color": SLATE,
         "tracking": 0.8, "caps": True})]], anchor=MSO_ANCHOR.MIDDLE)
    uw = track_w / units
    for u in range(units):
        text(slide, track_x + u*uw, cy, uw, 0.24, [[
            (f"W{u+1}", {"size": 7.5, "bold": True, "color": SLATE})]],
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    hline(slide, MARGIN_X, cy+0.26, CONTENT_W, GRAPHITE, 1.2)
    row_y = cy + 0.34
    row_h = (ch - 0.34) / len(bars)
    for label, ckey, start, span, note in bars:
        text(slide, MARGIN_X, row_y, label_w-0.1, row_h, [[
            (label, {"size": 10.5, "bold": True, "color": INK})]],
            anchor=MSO_ANCHOR.MIDDLE)
        bx = track_x + start*uw
        bw = span*uw
        bar_h = 0.32
        by = row_y + (row_h - bar_h)/2
        on_y = ckey == "yellow"
        rect(slide, bx, by, bw, bar_h, fill=MS[ckey]["base"], rounded=True, radius=0.16)
        text(slide, bx+0.08, by, bw-0.16, bar_h, [[
            (note, {"size": 8, "bold": True,
             "color": INK if on_y else WHITE})]], anchor=MSO_ANCHOR.MIDDLE)
        hline(slide, MARGIN_X, row_y+row_h, CONTENT_W, LINE2, 0.5)
        row_y += row_h

def r_decision_matrix(prs, s):
    slide, accent, ctop = content_frame(prs, s["sec"], s["kicker"],
        runs(s["h"], size=22, color=INK, bold_color=INK),
        runs(s["desc"], size=12.5, color=SLATE),
        runs(s["src"], size=10.5, color=SLATE, bold_color=GRAPHITE),
        title_size=22)
    headers = s["headers"]
    rows = s["rows"]
    total_row = s["total"]
    ncols = len(headers)
    crit_w = 0.34 * CONTENT_W
    num_w = (CONTENT_W - crit_w)/(ncols-1)
    col_w = [crit_w] + [num_w]*(ncols-1)
    ty = ctop
    nrows = len(rows) + 2
    HEADER_H = 0.42
    body_rh = min(0.46, (CONTENT_BOT - ty - HEADER_H - 1.0)/(len(rows)+1))
    tbl = native_table(slide, MARGIN_X, ty, CONTENT_W, nrows, ncols,
                       col_widths=col_w, row_h=0.40)
    tbl.rows[0].height = IN(HEADER_H)
    for ri in range(1, nrows):
        tbl.rows[ri].height = IN(body_rh)
    for ci, (label, is_num) in enumerate(headers):
        cell_text(tbl.cell(0,ci), [(label, {})], size=9, bold=True, color=SLATE,
                  align=PP_ALIGN.RIGHT if is_num else PP_ALIGN.LEFT,
                  fill=WHITE, anchor=MSO_ANCHOR.BOTTOM)
    for ri, (crit, cells) in enumerate(rows, start=1):
        zebra = BG_SOFT if ri % 2 == 0 else WHITE
        cell_text(tbl.cell(ri,0), [(crit, {"bold": True, "color": INK})], size=11, fill=zebra)
        for ci, val in enumerate(cells, start=1):
            cell_text(tbl.cell(ri,ci), [(val, {"bold": True})], size=11,
                      color=GRAPHITE, align=PP_ALIGN.RIGHT, fill=zebra)
    tr = len(rows)+1
    cell_text(tbl.cell(tr,0), [(total_row[0], {"bold": True, "color": INK})],
              size=11.5, fill=accent["l50"])
    for ci, val in enumerate(total_row[1], start=1):
        cell_text(tbl.cell(tr,ci), [(val, {"bold": True, "color": INK})],
                  size=12, align=PP_ALIGN.RIGHT, fill=accent["l50"])
    hline(slide, MARGIN_X, ty+HEADER_H, CONTENT_W, GRAPHITE, 1.5)
    hline(slide, MARGIN_X, ty+HEADER_H+body_rh*len(rows), CONTENT_W, GRAPHITE, 1.2)
    if s.get("pull"):
        py = ty + HEADER_H + body_rh*(len(rows)+1) + 0.18
        ph = CONTENT_BOT - py
        if ph > 0.55:
            lbl, txt, cite = s["pull"]
            rect(slide, MARGIN_X, py, CONTENT_W, ph, fill=accent["l50"], rounded=True, radius=0.04)
            rect(slide, MARGIN_X, py, 0.045, ph, fill=accent["base"])
            text(slide, MARGIN_X+0.24, py+0.12, CONTENT_W-0.5, 0.18, [[
                (lbl, {"size": 9, "bold": True, "color": accent["d700"],
                 "tracking": 1.8, "caps": True})]])
            text(slide, MARGIN_X+0.24, py+0.33, CONTENT_W-0.6, ph-0.5, [
                runs(txt, size=12, color=INK, bold_color=INK)
            ], line_spacing=1.3)

def r_stakeholder_map(prs, s):
    slide, accent, ctop = content_frame(prs, s["sec"], s["kicker"],
        runs(s["h"], size=22, color=INK, bold_color=INK),
        runs(s["desc"], size=12.5, color=SLATE),
        runs(s["src"], size=10.5, color=SLATE, bold_color=GRAPHITE),
        title_size=22)
    cy = ctop + 0.10
    yaxis_w = 0.34
    xaxis_h = 0.26
    qx = MARGIN_X + yaxis_w
    qw = CONTENT_W - yaxis_w
    grid_h = CONTENT_BOT - cy - xaxis_h - 0.10
    _vlabel(slide, MARGIN_X, cy, grid_h, s["yaxis"])
    cells = s["cells"]
    half_w = qw/2
    half_h = grid_h/2
    for idx, (label, chips, key) in enumerate(cells):
        r_, c_ = idx//2, idx%2
        cellx = qx + c_*half_w
        celly = cy + r_*half_h
        rect(slide, cellx+0.01, celly+0.01, half_w-0.02, half_h-0.02,
             fill=accent["l50"] if key else BG_CARD, line=LINE, line_w=1.0)
        text(slide, cellx+0.18, celly+0.14, half_w-0.36, 0.30, [[
            (label, {"size": 8, "bold": True,
             "color": accent["d700"] if key else SILVER,
             "tracking": 0.8, "caps": True})]])
        chx, chy = cellx+0.18, celly+0.48
        for chip in chips:
            cwid = 0.26 + len(chip)*0.072
            if chx + cwid > cellx + half_w - 0.18:
                chx = cellx+0.18
                chy += 0.34
            rect(slide, chx, chy, cwid, 0.26, fill=WHITE, line=LINE, line_w=1.0,
                 rounded=True, radius=0.5)
            text(slide, chx, chy, cwid, 0.26, [[
                (chip, {"size": 9.5, "color": GRAPHITE})]],
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            chx += cwid + 0.08
    text(slide, qx, cy+grid_h+0.04, qw, xaxis_h, [[
        (s["xaxis"], {"size": 8.5, "bold": True, "color": SLATE,
         "tracking": 1.0, "caps": True})]], align=PP_ALIGN.CENTER)

def r_org_chart(prs, s):
    slide, accent, ctop = content_frame(prs, s["sec"], s["kicker"],
        runs(s["h"], size=22, color=INK, bold_color=INK),
        runs(s["desc"], size=12.5, color=SLATE),
        runs(s["src"], size=10.5, color=SLATE, bold_color=GRAPHITE),
        title_size=22)
    cy = ctop + 0.10
    ch = CONTENT_BOT - cy - 0.05
    lead = s["lead"]
    reports = s["reports"]
    lead_w, lead_h = 2.6, 0.78
    lead_x = MARGIN_X + (CONTENT_W - lead_w)/2
    lead_y = cy + 0.30
    rect(slide, lead_x, lead_y, lead_w, lead_h, fill=BG_CARD, line=LINE, line_w=1.0,
         rounded=True, radius=0.10)
    rect(slide, lead_x, lead_y, lead_w, 0.05, fill=accent["base"])
    text(slide, lead_x, lead_y+0.14, lead_w, 0.28, [[
        (lead[0], {"size": 13, "bold": True, "color": INK})]], align=PP_ALIGN.CENTER)
    text(slide, lead_x, lead_y+0.42, lead_w, 0.24, [[
        (lead[1], {"size": 9.5, "color": SLATE})]], align=PP_ALIGN.CENTER)
    mid_x = MARGIN_X + CONTENT_W/2
    vline(slide, mid_x, lead_y+lead_h, 0.32, LINE, 1.5)
    rail_y = lead_y+lead_h+0.32
    n = len(reports)
    rgap = 0.30
    rw = (CONTENT_W*0.84 - rgap*(n-1))/n
    rail_start = MARGIN_X + (CONTENT_W - (rw*n + rgap*(n-1)))/2
    if n > 1:
        hline(slide, rail_start+rw/2, rail_y, (rw+rgap)*(n-1), LINE, 1.5)
    rh = 0.92
    for i, (title, sub) in enumerate(reports):
        rx = rail_start + i*(rw+rgap)
        vline(slide, rx+rw/2, rail_y, 0.26, LINE, 1.5)
        ryy = rail_y + 0.26
        rect(slide, rx, ryy, rw, rh, fill=BG_CARD, line=LINE, line_w=1.0,
             rounded=True, radius=0.08)
        text(slide, rx+0.14, ryy+0.16, rw-0.28, 0.4, [
            runs(title, size=12, color=INK, bold_color=INK)
        ], line_spacing=1.15, align=PP_ALIGN.CENTER)
        text(slide, rx+0.14, ryy+rh-0.36, rw-0.28, 0.3, [
            runs(sub, size=9, color=SLATE)
        ], line_spacing=1.25, align=PP_ALIGN.CENTER)

def r_appendix_divider(prs, s):
    slide = base_slide(prs, bg=INK)
    seg = PW/4
    for i, key in enumerate(["red","green","yellow","blue"]):
        rect(slide, i*seg, 0, seg, 0.06, fill=MS[key]["base"])
    text(slide, 1.1, 2.35, PW-2.2, 0.3, [[
        ("Appendix", {"size": 11, "bold": True, "color": MS["blue"]["base"],
         "tracking": 2.6, "caps": True})]])
    text(slide, 1.1, 2.70, PW-2.2, 0.9, [[
        ("Reference material", {"size": 44, "bold": True, "color": WHITE})]])
    text(slide, 1.1, 3.70, PW-3.5, 0.9, [
        runs("Supporting detail for the brief: the full assessment checklist, "
             "the service inventory, and the data behind each illustrative "
             "figure. Everything past this point is reference, not narrative.",
             size=14, color=RGBColor(0xBE,0xBE,0xBE))
    ], line_spacing=1.55)
    chips = ["Assessment checklist", "Service inventory", "Figure sources", "Glossary"]
    cx = 1.1
    for chip in chips:
        cw = 0.4 + len(chip)*0.082
        rect(slide, cx, 4.85, cw, 0.38, fill=None, line=RGBColor(0x44,0x44,0x44),
             line_w=1.0, rounded=True, radius=0.5)
        text(slide, cx, 4.85, cw, 0.38, [[
            (chip, {"size": 10.5, "color": RGBColor(0xD0,0xD0,0xD0)})]],
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        cx += cw + 0.14

def r_contact(prs, s):
    slide = base_slide(prs)
    half = PW/2
    rect(slide, 0, 0, half, PH, fill=MS["blue"]["base"])
    text(slide, 0.85, 2.00, half-1.5, 0.28, [[
        ("Next step", {"size": 11, "bold": True, "color": RGBColor(0xD8,0xEF,0xFB),
         "tracking": 2.4, "caps": True})]])
    text(slide, 0.85, 2.34, half-1.5, 1.7, [
        runs("Confirm an assessment window, and we begin.",
             size=29, color=WHITE, bold_color=WHITE)
    ], line_spacing=1.18)
    text(slide, 0.85, 4.55, half-1.7, 1.3, [
        runs("The three-week assessment is low-commitment and produces the "
             "client's own baseline. Everything after it is a decision made "
             "with real numbers.", size=13, color=RGBColor(0xE8,0xF6,0xFD))
    ], line_spacing=1.5)
    rx = half + 0.85
    rw = half - 1.7
    text(slide, rx, 1.85, rw, 0.20, [[
        ("What we need from you", {"size": 9.5, "bold": True,
         "color": MS["blue"]["d700"], "tracking": 1.6, "caps": True})]])
    text(slide, rx, 2.09, rw, 0.7, [
        runs("A **three-week window** and the two candidate workloads to baseline.",
             size=15, color=INK, bold_color=INK)
    ], line_spacing=1.4)
    text(slide, rx, 3.05, rw, 0.20, [[
        ("What you get back", {"size": 9.5, "bold": True,
         "color": MS["blue"]["d700"], "tracking": 1.6, "caps": True})]])
    text(slide, rx, 3.29, rw, 0.7, [
        runs("A measured baseline and a prioritized platform scope.",
             size=15, color=INK)
    ], line_spacing=1.4)
    hline(slide, rx, 4.40, rw, LINE, 0.75)
    text(slide, rx, 4.62, rw, 0.32, [[
        ("Frontier Cockpit Team", {"size": 16, "bold": True, "color": INK})]])
    text(slide, rx, 4.96, rw, 0.24, [[
        ("Software Global Black Belt", {"size": 11, "color": SLATE})]])
    text(slide, rx, 5.24, rw, 0.26, [[
        ("frontier-cockpit@example.com", {"size": 12, "bold": True,
         "color": MS["blue"]["d700"]})]])


RENDERERS = {
    "cover_meeting":       r_cover_meeting,
    "cover_stat":          r_cover_stat,
    "cover_thesis":        r_cover_thesis,
    "cover_pillars":       r_cover_pillars,
    "divider":             r_divider,
    "big_number":          r_big_number,
    "data_table":          r_data_table,
    "exec_summary":        r_exec_summary,
    "quote":               r_quote,
    "section_intro":       r_section_intro,
    "pillar_grid":         r_pillar_grid,
    "flow":                r_flow,
    "layer_stack":         r_layer_stack,
    "before_after":        r_before_after,
    "roadmap":             r_roadmap,
    "pricing":             r_pricing,
    "risks_table":         r_risks_table,
    "quadrant":            r_quadrant,
    "agenda":              r_agenda,
    "recommendation":      r_recommendation,
    "team":                r_team,
    "faq":                 r_faq,
    "metrics_dashboard":   r_metrics_dashboard,
    "case_study":          r_case_study,
    "architecture_detail": r_architecture_detail,
    "vendor_compare":      r_vendor_compare,
    "raci":                r_raci,
    "gantt":               r_gantt,
    "decision_matrix":     r_decision_matrix,
    "stakeholder_map":     r_stakeholder_map,
    "org_chart":           r_org_chart,
    "appendix_divider":    r_appendix_divider,
    "contact":             r_contact,
}

# ---- TEST DECK (partial, first batch of renderers) ----
TEST_DECK = [
    {"type": "cover_meeting", "sec": "front"},
    {"type": "divider", "sec": "prob", "act_no": "II",
     "kicker": "Section 2 \u00b7 The problem",
     "title": "Why pilots stall before production",
     "desc": "The blocker is rarely the model. It is everything underneath: "
             "environments, data access, governance, and the route to a service.",
     "questions": ["Where exactly do pilots break down?",
                   "What does the delay cost?",
                   "Is this a tooling problem or a platform problem?"],
     "q_start": 4},
    {"type": "big_number", "sec": "prob", "kicker": "The gap",
     "h": "Most teams reach a **working demo**. Far fewer reach a governed service.",
     "desc": "The distance between those two states is the platform gap.",
     "metric_label": "The core metric", "metric_n": "68", "metric_unit": "%",
     "metric_desc": "of AI pilots in this scenario never get a production target assigned.",
     "cards": [("Symptom 1 \u00b7 Environments", "~40h", "spent per team, per quarter, re-creating environments by hand.", ""),
               ("Symptom 2 \u00b7 Governance", "0", "consistent policy gates between an experiment and an endpoint.", ""),
               ("Root cause", "1", "missing platform layer. Every symptom traces back to it.", "accent")],
     "src": "**Note \u00b7** All figures are placeholders for the reference template."},
    {"type": "data_table", "sec": "prob", "kicker": "Cost of delay",
     "h": "The platform gap has a **line-item cost**. Here is where it accrues.",
     "desc": "Five recurring friction points, the team they tax, and the effect.",
     "headers": [("Friction point", "l", 0.24), ("Owner", "l", 0.16),
                 ("Effect on delivery", "l", 0.30), ("Recurs", "num", 0.13),
                 ("Severity", "l", 0.17)],
     "rows": [
        [("Manual environment setup", "bold"), ("Platform team", None),
         ("Each new workload restarts from zero infrastructure.", None),
         ("Per project", "num"), ("High", "tag")],
        [("Inconsistent data access", "bold"), ("Data engineering", None),
         ("Credentials and patterns differ across every team.", None),
         ("Per team", "num"), ("High", "tag")],
        [("No standard CI path", "bold"), ("App teams", None),
         ("Each repo invents its own build and release flow.", None),
         ("Per repo", "num"), ("Medium", "tag")],
        [("Governance applied late", "bold"), ("Security", None),
         ("Policy checks happen after build, forcing rework.", None),
         ("Per release", "num"), ("Medium", "tag")],
        [("Fragmented observability", "bold"), ("SRE", None),
         ("No shared view of cost, latency, or model behavior.", None),
         ("Continuous", "num"), ("Watch", "tag")],
     ],
     "pull": ("Read this as",
              "Every row is **solvable once** at the platform layer instead of "
              "**repeatedly** at the project layer.",
              "Synthesis line \u00b7 reference template"),
     "src": "**Severity banding** is illustrative. Replace with the client's own data."},
    {"type": "exec_summary", "sec": "front", "kicker": "Executive summary",
     "h": "Three things to take from this brief, **before the detail**.",
     "desc": "If the room reads only one slide, it is this one.",
     "cols": [
        ("1", "The blocker is the platform, not the model",
         "Pilots reach a working demo and stall. The recurring cause is the absence of a shared foundation: environments, data access, and a route to production.",
         "**So what \u00b7** Solve it once at the platform layer, not per project."),
        ("2", "An opinionated stack closes the gap",
         "Four layers, each owned once, remove the friction line-items teams hit today. Teams consume the platform; they do not re-derive it.",
         "**So what \u00b7** The reference architecture is the deliverable."),
        ("3", "Start small, decide on evidence",
         "A three-week assessment produces the client's own baseline. Every commitment after that is a decision made with real numbers.",
         "**So what \u00b7** Low-commitment first step. The roadmap earns the next."),
     ],
     "src": "**Read order \u00b7** This slide stands alone; the six sections expand each column."},
    {"type": "quote", "sec": "prob",
     "quote": "We did not have an AI problem. We had twelve different "
              "definitions of what it meant to put a model in production.",
     "cite": "Practitioner voice \u00b7 **reference template**, illustrative attribution"},
    {"type": "section_intro", "sec": "ctx", "kicker": "Section intro",
     "h": "Where the market is, in **one honest read**.",
     "desc": "A section intro is lighter than a divider; it frames the sub-topic.",
     "eyebrow": "What this section covers",
     "statement": "Adoption is wide. **Production is narrow.**",
     "lead": "Most enterprises now run AI pilots. Far fewer run AI in production. "
             "This section sizes that gap with the client's peers as the "
             "reference point, so the problem is concrete before we name its causes.",
     "aside": [("In this section", "4 slides, ending on a client snapshot"),
               ("Key question", "How big is the pilot-to-production gap?"),
               ("Leaves you with", "A measured starting point, not an opinion")],
     "src": "**Pattern note \u00b7** Use a section intro after a divider when the section is long."},
    {"type": "pillar_grid", "sec": "prob", "kicker": "Root causes",
     "h": "Five symptoms, **three root causes**.",
     "desc": "The cost-of-delay table lists what teams feel. These are the causes.",
     "pillars": [
        ("01", "No shared foundation", "Every workload re-derives networking, identity, and a release path because nothing provides them by default."),
        ("02", "Governance bolted on", "Policy lives outside the delivery path, so it is discovered late and applied as rework rather than as a property of shipping."),
        ("03", "No common signal", "Cost, latency, and model behavior live in team-local dashboards, so problems surface in weeks instead of hours."),
     ],
     "src": "**Cause note \u00b7** Three is deliberate. A long list of causes is a list nobody acts on."},
    {"type": "flow", "sec": "sol", "kicker": "The governed path",
     "h": "With the platform in place, a workload moves through **four governed stages**.",
     "desc": "Each stage has a clear owner and a clear hand-off.",
     "nodes": [
        ("1", "Scaffold", "Start from a golden-path template with identity and networking pre-wired."),
        ("2", "Build", "Standard CI runs tests, evaluations, and policy checks in one pass."),
        ("3", "Deploy", "Release to the AKS landing zone with scaling and rollback defined."),
        ("4", "Operate", "Shared observability tracks cost, latency, and model behavior."),
     ],
     "src": "**Flow note \u00b7** This is the broken-path slide from Section 2, fixed stage for stage."},
]


def build(deck, only=None):
    prs = Presentation()
    prs.slide_width = IN(PW)
    prs.slide_height = IN(PH)
    CUR["total"] = len(deck)
    rendered = 0
    for i, spec in enumerate(deck, 1):
        if only and spec["type"] not in only:
            continue
        CUR["page"] = i
        fn = RENDERERS.get(spec["type"])
        if fn is None:
            sys.exit(f"ERROR: no renderer for '{spec['type']}' (slide {i})")
        before = len(prs.slides)
        fn(prs, spec)
        if len(prs.slides) <= before:
            sys.exit(f"ERROR: renderer '{spec['type']}' did not create a slide (slide {i})")
        attach_speaker_note(prs.slides[len(prs.slides) - 1], spec, i, CUR["total"])
        rendered += 1
    if rendered == 0:
        sys.exit("ERROR: no slides rendered; check --only or the DECK list")
    return prs


def verify_output(path, expected_slides):
    out = Path(path)
    if not out.is_file() or out.stat().st_size == 0:
        sys.exit(f"ERROR: PPTX output is missing or empty: {out}")
    try:
        check = Presentation(str(out))
    except Exception as exc:
        sys.exit(f"ERROR: PPTX output cannot be reopened: {exc}")
    actual = len(check.slides)
    if actual != expected_slides:
        sys.exit(f"ERROR: PPTX has {actual} slides, expected {expected_slides}")
    missing_notes = []
    for index, slide in enumerate(check.slides, start=1):
        note_text = slide.notes_slide.notes_text_frame.text if slide.has_notes_slide else ""
        if not note_text.strip() or "[OPENING]" not in note_text or "[CORE]" not in note_text:
            missing_notes.append(index)
    if missing_notes:
        sys.exit(f"ERROR: PPTX speaker notes missing or incomplete on slides {missing_notes[:12]}")
    return out.stat().st_size, actual


if __name__ == "__main__":
    ap = argparse.ArgumentParser()
    ap.add_argument("--out", default="ms-reference-deck.pptx")
    ap.add_argument("--only", default=None, help="comma-separated slide types")
    ap.add_argument("--test", action="store_true", help="build the partial TEST_DECK")
    args = ap.parse_args()
    if args.test:
        deck = TEST_DECK
    else:
        from deck_data import DECK
        deck = DECK
    only = set(args.only.split(",")) if args.only else None
    prs = build(deck, only)
    prs.save(args.out)
    expected = len([s for s in deck if not only or s["type"] in only])
    size, n = verify_output(args.out, expected)
    print(f"saved {args.out} ({n} slides, {size} bytes)")
