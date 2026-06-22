#!/usr/bin/env python3
"""Initialize and validate GitHub Copilot UBB client package folders."""

from __future__ import annotations

import argparse
import hashlib
import json
import re
import shutil
import subprocess
import sys
import tempfile
import unicodedata
from collections import Counter
from datetime import datetime, timezone
from pathlib import Path


DEFAULT_OUTPUT_ROOT = Path("client-ubb-transition-kits")
DEFAULT_REFERENCE_ROOT = Path("REFERENCES_gh-ubb-customers-kit")
DEFAULT_TRANSITION_PLAN_SOURCE = DEFAULT_REFERENCE_ROOT / \
    "00-Customers_Transition_Plan"
DB_FILENAME = "db.json"
MANIFEST_FILENAME = "package-manifest.json"
DECK_DERIVATIVE_LOCALES = ("pt-BR", "en", "es")
DEFAULT_PDF_LOCALES = ("pt-BR",)
PDF_VISUAL_QA_DPI = 120
SOFFICE_TIMEOUT_SECONDS = 180

TRANSITION_PLAN_TPIDS = {
    "B3": 2011058,
    "BTG": 5336866,
    "BancoDoBrasil": 524098,
    "Boticario": 3287906,
    "Bradesco": 6919482,
    "Energisa": 1523937,
    "Globo": 1427980,
    "GrupoFolha": 22513842,
    "Intermedica": 955113,
    "Itau": 1719150,
    "Localiza": 967068,
    "MagazineLuiza": 15606168,
    "PETROBRAS": 940486,
    "SERPRO": 13088972,
    "Sicredi": 10015354,
    "Stone": 23102848,
    "XP": 16751888,
}

REQUIRED_SECTIONS = (
    "ranking_brasil",
    "clientes",
    "cliente_modelo",
    "modelos_tokens",
    "uso_mensal",
    "billable_owners",
)

CRITICAL_RECONCILIATION_FIELDS = (
    "top_parent",
    "ubb_gasto_atual_licenca_pru_usd",
    "ubb_promo_gasto_estimado_usd",
    "ubb_standard_gasto_estimado_usd",
    "ubb_promo_diff_usd",
    "ubb_standard_diff_usd",
    "ghcp_seats",
    "ghcp_acr_usd",
    "pru_acr_usd",
)

GENERIC_PDFS = (
    "Runbook Otimizacao GitHubCopilot UBB v4 0 0 2026-06-09 ptBR.pdf",
    "Playbook Implementacao GitHubCopilot UBB v4 0 0 2026-06-09 ptBR.pdf",
    "GitHubCopilot_UBB_Runbook_TechDoc_v5_0_0_2026-06-11_ptBR.pdf",
    "GitHubCopilot_UBB_Otimizacao_Guia_v3_2_0_2026-06-10_ptBR.pdf",
)

PACKAGE_DIRS = (
    "INTERNAL_Budget_Impact_Analysis",
    "INTERNAL_Transition_Kit",
    "EXTERNAL_Customer_Ready_Documents/customer-deck",
    "EXTERNAL_Customer_Ready_Documents/customer-ready-playbooks",
)

CUSTOM_DELIVERABLES = {
    "internal_budget_html": (
        Path("INTERNAL_Budget_Impact_Analysis"),
        re.compile(
            r"^(?:INTERNAL_ONLY|USO INTERNO) - .+Analise_Budget.+_multi\.html$"),
    ),
    "internal_transition_xlsx": (
        Path("INTERNAL_Transition_Kit"),
        re.compile(
            r"^(?:INTERNAL_ONLY|USO INTERNO) - Kit_Transicao_UBB_.+_EN\.xlsx$"),
    ),
    "customer_deck_html": (
        Path("EXTERNAL_Customer_Ready_Documents/customer-deck"),
        re.compile(
            r"^(?:CUSTOMER_Deck_.+_06_Mo(?:n|un)th_Program.+|.+_Programa_Seis_Meses_Deck.+)_multi\.html$"),
    ),
    "customer_proposal_pdf": (
        Path("EXTERNAL_Customer_Ready_Documents/customer-ready-playbooks"),
        re.compile(r"^Proposta_Acao_UBB_.+_Editorial_.+_ptBR\.pdf$"),
    ),
}

SOURCE_CUSTOM_DELIVERABLES = {
    "internal_budget_html": re.compile(r"^USO INTERNO - .+Analise_Budget.+_multi\.html$"),
    "internal_transition_xlsx": re.compile(r"^USO INTERNO - Kit_Transicao_UBB_.+_EN\.xlsx$"),
    "customer_deck_html": re.compile(r"^.+_Programa_Seis_Meses_Deck.+_multi\.html$"),
    "customer_proposal_pdf": re.compile(r"^Proposta_Acao_UBB_.+_Editorial_.+_ptBR\.pdf$"),
}


class Reporter:
    def __init__(self) -> None:
        self.errors: list[str] = []
        self.source_gaps: list[str] = []
        self.warnings: list[str] = []

    def error(self, message: str) -> None:
        self.errors.append(message)
        print(f"FAIL {message}")

    def warn(self, message: str) -> None:
        self.warnings.append(message)
        print(f"WARN {message}")

    def source_gap(self, message: str) -> None:
        self.source_gaps.append(message)
        print(f"SOURCE-GAP {message}")

    def pass_(self, message: str) -> None:
        print(f"PASS {message}")

    def exit_code(self, allow_source_gaps: bool = False) -> int:
        if self.errors:
            return 1
        if self.source_gaps and not allow_source_gaps:
            return 1
        return 0


def load_json(path: Path) -> dict:
    try:
        return json.loads(path.read_text(encoding="utf-8"))
    except FileNotFoundError:
        raise SystemExit(f"Missing db file: {path}")
    except json.JSONDecodeError as exc:
        raise SystemExit(f"Invalid JSON in {path}: {exc}")


def client_slug(name: str, tpid: int) -> str:
    normalized = unicodedata.normalize("NFKD", name)
    ascii_name = normalized.encode("ascii", "ignore").decode("ascii")
    slug = re.sub(r"[^A-Za-z0-9]+", "_", ascii_name).strip("_")
    slug = re.sub(r"_+", "_", slug)
    return f"{slug}_{tpid}"


def tpid_key(row: dict) -> int:
    return int(row["tpid"])


def rows_for(data: dict, section: str, tpid: int) -> list[dict]:
    return [row for row in data.get(section, []) if tpid_key(row) == tpid]


def single_row(data: dict, section: str, tpid: int) -> dict | None:
    rows = rows_for(data, section, tpid)
    if not rows:
        return None
    return rows[0]


def ranked_tpids(data: dict) -> set[int]:
    return {tpid_key(row) for row in data.get("ranking_brasil", [])}


def all_section_tpids(data: dict, section: str) -> set[int]:
    return {tpid_key(row) for row in data.get(section, []) if "tpid" in row}


def compare_values(left: object, right: object) -> bool:
    if isinstance(left, (int, float)) and isinstance(right, (int, float)):
        return abs(float(left) - float(right)) <= 0.02
    return left == right


def validate_top_level_keys(data: dict, reporter: Reporter) -> bool:
    missing_top = [key for key in (
        "meta", "ranking_brasil") if key not in data]
    for key in missing_top:
        reporter.error(f"db.json missing top-level key: {key}")
    return not missing_top


def validate_required_sections(data: dict, tpids: set[int], reporter: Reporter) -> None:
    for section in REQUIRED_SECTIONS:
        if section not in data or not isinstance(data[section], list):
            reporter.error(f"db.json missing required list section: {section}")
            continue
        have = all_section_tpids(data, section)
        missing = sorted(tpids - have)
        extra = sorted(have - tpids)
        if missing:
            names = {tpid_key(row): row.get("top_parent", "unknown")
                     for row in data["ranking_brasil"]}
            detail = ", ".join(
                f"{tpid} ({names.get(tpid, 'unknown')})" for tpid in missing[:10])
            reporter.source_gap(
                f"section {section} missing TPID rows: {detail}")
        else:
            reporter.pass_(f"section {section} covers all ranked clients")
        if extra:
            reporter.warn(
                f"section {section} has TPIDs outside ranking_brasil: {extra[:10]}")


def validate_reconciliation(data: dict, reporter: Reporter) -> None:
    clients_by_tpid = {tpid_key(row): row for row in data.get("clientes", [])}
    for ranking in data.get("ranking_brasil", []):
        tpid = tpid_key(ranking)
        client = clients_by_tpid.get(tpid)
        if not client:
            continue
        for field in CRITICAL_RECONCILIATION_FIELDS:
            if field not in ranking or field not in client:
                reporter.error(
                    f"TPID {tpid} missing reconciliation field {field}")
                continue
            if not compare_values(ranking[field], client[field]):
                reporter.error(
                    f"TPID {tpid} reconciliation mismatch for {field}: "
                    f"ranking_brasil={ranking[field]} clientes={client[field]}"
                )


def warn_partial_monthly_history(data: dict, reporter: Reporter) -> None:
    monthly_counts = Counter(tpid_key(row)
                             for row in data.get("uso_mensal", []))
    for tpid, count in sorted(monthly_counts.items()):
        if count < 11:
            name = single_row(data, "ranking_brasil", tpid).get(
                "top_parent", "unknown")
            reporter.warn(
                f"TPID {tpid} ({name}) has {count} monthly usage rows, less than the 11-row full FY26 reference")


def validate_db(data: dict, reporter: Reporter) -> None:
    if not validate_top_level_keys(data, reporter):
        return

    tpids = ranked_tpids(data)
    reporter.pass_(f"ranking_brasil contains {len(tpids)} clients")
    validate_required_sections(data, tpids, reporter)
    validate_reconciliation(data, reporter)
    warn_partial_monthly_history(data, reporter)


def source_gaps(data: dict, tpid: int) -> list[str]:
    gaps: list[str] = []
    for section in REQUIRED_SECTIONS:
        if not rows_for(data, section, tpid):
            gaps.append(section)
    return gaps


def client_snapshot(data: dict, source_path: Path, tpid: int) -> dict:
    client = single_row(data, "clientes", tpid)
    ranking = single_row(data, "ranking_brasil", tpid)
    if not client or not ranking:
        raise SystemExit(
            f"TPID {tpid} is not present in clientes and ranking_brasil")
    return {
        "source": {
            "file": str(source_path),
            "version": data.get("meta", {}).get("versao"),
            "generated_at": data.get("meta", {}).get("gerado_em"),
        },
        "meta": data.get("meta", {}),
        "client": client,
        "ranking_brasil": ranking,
        "cliente_modelo": rows_for(data, "cliente_modelo", tpid),
        "modelos_tokens": rows_for(data, "modelos_tokens", tpid),
        "uso_mensal": rows_for(data, "uso_mensal", tpid),
        "billable_owners": rows_for(data, "billable_owners", tpid),
    }


def package_name(data: dict, tpid: int) -> str:
    ranking = single_row(data, "ranking_brasil", tpid)
    if not ranking:
        raise SystemExit(f"TPID {tpid} is not present in ranking_brasil")
    return client_slug(str(ranking["top_parent"]), tpid)


def generic_pdf_source_dir(reference_root: Path) -> Path:
    return reference_root / "EXTERNAL_Customer_Ready_Documents" / "customer-ready-playbooks"


def sha256_file(path: Path) -> str:
    digest = hashlib.sha256()
    with path.open("rb") as handle:
        for chunk in iter(lambda: handle.read(1024 * 1024), b""):
            digest.update(chunk)
    return digest.hexdigest()


def artifact_record(category: str, kind: str, source_path: Path, destination_path: Path) -> dict:
    return {
        "category": category,
        "kind": kind,
        "source_path": str(source_path),
        "destination_path": str(destination_path),
        "sha256": sha256_file(destination_path),
        "bytes": destination_path.stat().st_size,
    }


def transition_source_folders(source_root: Path, selected: list[str] | None = None) -> list[Path]:
    selected_set = set(selected or [])
    folders = [path for path in source_root.iterdir() if path.is_dir()
               and not path.name.startswith(".")]
    if selected_set:
        folders = [path for path in folders if path.name in selected_set]
        missing = sorted(selected_set - {path.name for path in folders})
        if missing:
            raise SystemExit(
                f"Requested client folders not found: {', '.join(missing)}")
    return sorted(folders, key=lambda path: path.name.lower())


def resolve_transition_tpid(folder_name: str) -> int:
    if folder_name not in TRANSITION_PLAN_TPIDS:
        known = ", ".join(sorted(TRANSITION_PLAN_TPIDS))
        raise ValueError(
            f"No deterministic TPID mapping for folder {folder_name}. Known folders: {known}")
    return TRANSITION_PLAN_TPIDS[folder_name]


def find_source_deliverables(source_dir: Path, reporter: Reporter) -> dict[str, Path]:
    deliverables: dict[str, Path] = {}
    files = [path for path in source_dir.iterdir() if path.is_file()]
    for label, pattern in SOURCE_CUSTOM_DELIVERABLES.items():
        matches = sorted(path for path in files if pattern.match(path.name))
        if len(matches) == 1:
            deliverables[label] = matches[0]
            reporter.pass_(
                f"source {source_dir.name} has {label}: {matches[0].name}")
        elif len(matches) == 0:
            reporter.error(f"source {source_dir} missing {label}")
        else:
            names = ", ".join(path.name for path in matches)
            reporter.error(
                f"source {source_dir} has multiple {label} files: {names}")
    return deliverables


def move_to_archive(path: Path) -> None:
    archive_dir = path.parent / "archive"
    archive_dir.mkdir(parents=True, exist_ok=True)
    target = archive_dir / path.name
    if target.exists():
        stem = path.stem
        suffix = path.suffix
        counter = 1
        while target.exists():
            target = archive_dir / f"{stem}_{counter}{suffix}"
            counter += 1
    shutil.move(str(path), str(target))


def archive_existing_current_files(package_dir: Path, label: str, destination_name: str) -> None:
    relative_dir, pattern = CUSTOM_DELIVERABLES[label]
    destination_dir = package_dir / relative_dir
    for existing in non_archive_matches(destination_dir, pattern):
        if existing.name != destination_name:
            move_to_archive(existing)


def copy_custom_deliverables(source_dir: Path, package_dir: Path, reporter: Reporter) -> list[dict]:
    artifacts: list[dict] = []
    deliverables = find_source_deliverables(source_dir, reporter)
    if len(deliverables) != len(SOURCE_CUSTOM_DELIVERABLES):
        return artifacts
    for label, source_file in deliverables.items():
        relative_dir, _pattern = CUSTOM_DELIVERABLES[label]
        destination = package_dir / relative_dir / source_file.name
        archive_existing_current_files(package_dir, label, source_file.name)
        shutil.copy2(source_file, destination)
        reporter.pass_(f"copied {label}: {destination}")
        artifacts.append(artifact_record(
            label, "custom", source_file, destination))
    return artifacts


def generic_artifact_records(reference_root: Path, package_dir: Path) -> list[dict]:
    source_dir = generic_pdf_source_dir(reference_root)
    target_dir = package_dir / "EXTERNAL_Customer_Ready_Documents" / \
        "customer-ready-playbooks"
    records = []
    for name in GENERIC_PDFS:
        records.append(artifact_record("generic_pdf", "generic",
                       source_dir / name, target_dir / name))
    return records


def write_manifest(
    data: dict,
    db_path: Path,
    source_dir: Path,
    package_dir: Path,
    tpid: int,
    artifacts: list[dict],
    reporter: Reporter,
) -> None:
    ranking = single_row(data, "ranking_brasil", tpid) or {}
    manifest = {
        "schema_version": "1.0.0",
        "imported_at_utc": datetime.now(timezone.utc).replace(microsecond=0).isoformat(),
        "source_db": str(db_path),
        "source_db_version": data.get("meta", {}).get("versao"),
        "source_db_generated_at": data.get("meta", {}).get("gerado_em"),
        "source_folder": str(source_dir),
        "package_folder": str(package_dir),
        "tpid": tpid,
        "top_parent": ranking.get("top_parent"),
        "artifacts": artifacts,
    }
    manifest_path = package_dir / MANIFEST_FILENAME
    manifest_path.write_text(json.dumps(
        manifest, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")
    reporter.pass_(f"wrote package manifest: {manifest_path}")


def load_manifest(package_dir: Path) -> dict:
    manifest_path = package_dir / MANIFEST_FILENAME
    if not manifest_path.exists():
        raise FileNotFoundError(f"package manifest not found: {manifest_path}")
    return load_json(manifest_path)


def update_manifest_artifacts(package_dir: Path, artifacts: list[dict], reporter: Reporter) -> None:
    manifest_path = package_dir / MANIFEST_FILENAME
    manifest = load_manifest(package_dir)
    existing = manifest.get("artifacts", [])
    new_destinations = {artifact["destination_path"] for artifact in artifacts}
    manifest["artifacts"] = [artifact for artifact in existing if artifact.get(
        "destination_path") not in new_destinations]
    manifest["artifacts"].extend(artifacts)
    manifest["derivatives_updated_at_utc"] = datetime.now(
        timezone.utc).replace(microsecond=0).isoformat()
    manifest_path.write_text(json.dumps(
        manifest, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")
    reporter.pass_(f"updated package manifest: {manifest_path}")


def deck_base(html_path: Path) -> str:
    stem = html_path.stem
    return stem[:-6] if stem.endswith("_multi") else stem


def deck_slide_count(html_path: Path) -> int:
    html = html_path.read_text(encoding="utf-8")
    count = len(re.findall(
        r"<section\b(?=[^>]*\bclass=[\"'][^\"']*\bslide\b)", html, flags=re.I))
    if count == 0:
        raise ValueError(f"deck has no slides: {html_path}")
    return count


def customer_deck_html(package_dir: Path) -> Path:
    deck_dir = package_dir / "EXTERNAL_Customer_Ready_Documents" / "customer-deck"
    matches = sorted(path for path in deck_dir.glob(
        "*_multi.html") if path.is_file())
    if len(matches) != 1:
        raise ValueError(
            f"expected exactly one customer deck HTML in {deck_dir}, found {len(matches)}")
    return matches[0]


def run_checked(cmd: list[str], reporter: Reporter, timeout: int | None = None) -> None:
    try:
        result = subprocess.run(
            cmd, cwd=Path.cwd(), capture_output=True, text=True, timeout=timeout)
    except subprocess.TimeoutExpired as exc:
        raise RuntimeError(
            f"command timed out after {timeout}s: {' '.join(cmd)}") from exc
    if result.stdout:
        print(result.stdout, end="")
    if result.stderr:
        print(result.stderr, end="", file=sys.stderr)
    if result.returncode != 0:
        raise RuntimeError(
            f"command failed ({result.returncode}): {' '.join(cmd)}")
    reporter.pass_(f"command ok: {' '.join(cmd[:2])}")


def copy_derivative(canonical_path: Path, package_dir: Path, relative_dir: Path, html_path: Path, category: str) -> dict:
    target_dir = package_dir / "EXTERNAL_Customer_Ready_Documents" / \
        "customer-deck" / relative_dir / deck_base(html_path)
    target_dir.mkdir(parents=True, exist_ok=True)
    target = target_dir / canonical_path.name
    shutil.copy2(canonical_path, target)
    return artifact_record(category, "derivative", html_path, target)


def generate_deck_derivatives_for_package(
    package_dir: Path,
    pdf_locales: tuple[str, ...],
    pptx_locales: tuple[str, ...],
    reporter: Reporter,
) -> None:
    html_path = customer_deck_html(package_dir)
    base = deck_base(html_path)
    total = deck_slide_count(html_path)
    canonical_pdf_dir = Path("html") / "decks" / "pdf" / base
    canonical_pptx_dir = Path("html") / "decks" / "pptx" / base
    canonical_pdf_dir.mkdir(parents=True, exist_ok=True)
    canonical_pptx_dir.mkdir(parents=True, exist_ok=True)

    records: list[dict] = []
    for locale in pdf_locales:
        pdf_path = canonical_pdf_dir / f"{base}_{locale}.pdf"
        run_checked([
            "python3",
            ".github/skills/ms-presentation-deck/scripts/make_pdf.py",
            str(html_path),
            str(pdf_path),
            "--total",
            str(total),
            "--locale",
            locale,
        ], reporter)
        run_checked([
            "python3",
            ".github/skills/ms-presentation-deck/scripts/validate_derivatives.py",
            "--html",
            str(html_path),
            "--locale",
            locale,
            "--pdf",
            str(pdf_path),
        ], reporter)
        records.append(copy_derivative(pdf_path, package_dir,
                       Path("pdf"), html_path, "customer_deck_pdf"))

    for locale in pptx_locales:
        pptx_path = canonical_pptx_dir / f"{base}_{locale}.pptx"
        run_checked([
            "python3",
            ".github/skills/ms-presentation-deck/scripts/build_native_pptx.py",
            "--html",
            str(html_path),
            "--locale",
            locale,
            "--output",
            str(pptx_path),
        ], reporter)
        records.append(copy_derivative(pptx_path, package_dir,
                       Path("pptx"), html_path, "customer_deck_pptx"))

    update_manifest_artifacts(package_dir, records, reporter)


def require_python_module(module_name: str, install_hint: str) -> None:
    try:
        __import__(module_name)
    except ImportError as exc:
        raise RuntimeError(
            f"missing Python module {module_name}. {install_hint}") from exc


def render_pdf_to_images(pdf_path: Path, output_dir: Path, prefix: str, reporter: Reporter) -> list[Path]:
    pdftoppm = shutil.which("pdftoppm")
    if not pdftoppm:
        raise RuntimeError("pdftoppm is required for visual PDF QA")
    output_dir.mkdir(parents=True, exist_ok=True)
    output_prefix = output_dir / prefix
    run_checked([pdftoppm, "-jpeg", "-r", str(PDF_VISUAL_QA_DPI),
                str(pdf_path), str(output_prefix)], reporter)
    images = sorted(output_dir.glob(f"{prefix}-*.jpg"))
    if not images:
        raise RuntimeError(f"pdftoppm produced no page images for {pdf_path}")
    return images


def image_nonwhite_ratio(image_path: Path) -> float:
    from PIL import Image

    with Image.open(image_path) as image:
        gray = image.convert("L").resize((160, 90))
        histogram = gray.histogram()
        nonwhite = sum(histogram[:245])
        return nonwhite / float(gray.width * gray.height)


def make_contact_sheet(image_paths: list[Path], output_path: Path, title: str, columns: int = 5) -> None:
    from PIL import Image, ImageDraw

    if not image_paths:
        raise RuntimeError(f"cannot create empty contact sheet: {output_path}")
    thumb_w = 320
    label_h = 28
    margin = 18
    title_h = 42
    thumbs = []
    for image_path in image_paths:
        image = Image.open(image_path).convert("RGB")
        ratio = thumb_w / image.width
        thumb_h = int(image.height * ratio)
        thumbs.append((image.resize((thumb_w, thumb_h)), image_path.name))
    thumb_h = max(thumb.height for thumb, _name in thumbs)
    rows = (len(thumbs) + columns - 1) // columns
    sheet_w = margin + columns * (thumb_w + margin)
    sheet_h = title_h + margin + rows * (thumb_h + label_h + margin)
    sheet = Image.new("RGB", (sheet_w, sheet_h), "white")
    draw = ImageDraw.Draw(sheet)
    draw.text((margin, 12), title, fill=(30, 30, 30))
    for index, (thumb, name) in enumerate(thumbs, start=1):
        row = (index - 1) // columns
        col = (index - 1) % columns
        x = margin + col * (thumb_w + margin)
        y = title_h + margin + row * (thumb_h + label_h + margin)
        sheet.paste(thumb, (x, y))
        draw.rectangle((x, y, x + thumb_w - 1, y + thumb_h - 1),
                       outline=(210, 210, 210))
        draw.text((x, y + thumb_h + 6),
                  f"{index:02d}  {name}", fill=(60, 60, 60))
    output_path.parent.mkdir(parents=True, exist_ok=True)
    sheet.save(output_path, quality=88)


def qa_pdf(pdf_path: Path, expected_pages: int, qa_dir: Path, reporter: Reporter) -> dict:
    from pypdf import PdfReader

    page_dir = qa_dir / "pages" / pdf_path.stem
    if page_dir.exists():
        shutil.rmtree(page_dir)
    reader = PdfReader(str(pdf_path))
    page_count = len(reader.pages)
    text_pages = sum(1 for page in reader.pages if (
        page.extract_text() or "").strip())
    images = render_pdf_to_images(pdf_path, page_dir, "page", reporter)
    nonwhite = [image_nonwhite_ratio(path) for path in images]
    blank_pages = [index for index, ratio in enumerate(
        nonwhite, start=1) if ratio < 0.005]
    contact_sheet = qa_dir / f"{pdf_path.stem}_contact_sheet.jpg"
    make_contact_sheet(images, contact_sheet,
                       f"PDF visual QA: {pdf_path.name}")
    status = page_count == expected_pages and len(
        images) == expected_pages and text_pages == expected_pages and not blank_pages
    if status:
        reporter.pass_(f"PDF visual QA passed: {pdf_path}")
    else:
        reporter.error(f"PDF visual QA failed: {pdf_path}")
    return {
        "path": str(pdf_path),
        "pages": page_count,
        "rendered_pages": len(images),
        "text_pages": text_pages,
        "blank_pages": blank_pages,
        "contact_sheet": str(contact_sheet),
        "status": "PASS" if status else "FAIL",
    }


def pptx_text_contact_sheet(pptx_path: Path, qa_dir: Path, locale: str) -> Path:
    from PIL import Image, ImageDraw
    from pptx import Presentation

    prs = Presentation(str(pptx_path))
    columns = 5
    card_w = 330
    card_h = 170
    margin = 18
    title_h = 42
    rows = (len(prs.slides) + columns - 1) // columns
    sheet = Image.new("RGB", (margin + columns * (card_w + margin),
                      title_h + margin + rows * (card_h + margin)), "white")
    draw = ImageDraw.Draw(sheet)
    draw.text((margin, 12),
              f"PPTX structural QA: {pptx_path.name} ({locale})", fill=(30, 30, 30))
    for index, slide in enumerate(prs.slides, start=1):
        row = (index - 1) // columns
        col = (index - 1) % columns
        x = margin + col * (card_w + margin)
        y = title_h + margin + row * (card_h + margin)
        draw.rectangle((x, y, x + card_w, y + card_h),
                       outline=(210, 210, 210), fill=(250, 250, 248))
        texts = [shape.text.strip() for shape in slide.shapes if getattr(
            shape, "has_text_frame", False) and shape.text.strip()]
        preview = " ".join(texts)[:260]
        draw.text((x + 10, y + 10), f"Slide {index:02d}", fill=(0, 120, 180))
        draw.text((x + 10, y + 34), preview, fill=(40, 40, 40))
    output_path = qa_dir / f"{pptx_path.stem}_structure_contact_sheet.jpg"
    output_path.parent.mkdir(parents=True, exist_ok=True)
    sheet.save(output_path, quality=88)
    return output_path


def pptx_slide_stats(prs) -> dict:
    slides_with_text = 0
    slides_with_notes = 0
    full_slide_images = 0
    for slide in prs.slides:
        if any(getattr(shape, "has_text_frame", False) and shape.text.strip() for shape in slide.shapes):
            slides_with_text += 1
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame.text.strip():
            slides_with_notes += 1
        full_slide_images += count_full_slide_images(slide, prs)
    return {
        "slides_with_text": slides_with_text,
        "slides_with_notes": slides_with_notes,
        "full_slide_images": full_slide_images,
    }


def count_full_slide_images(slide, prs) -> int:
    total = 0
    for shape in slide.shapes:
        is_picture = getattr(shape, "shape_type", None) == 13
        is_full_width = shape.width >= prs.slide_width * 0.95
        is_full_height = shape.height >= prs.slide_height * 0.95
        if is_picture and is_full_width and is_full_height:
            total += 1
    return total


def qa_pptx(pptx_path: Path, expected_slides: int, qa_dir: Path, reporter: Reporter) -> dict:
    from pptx import Presentation

    prs = Presentation(str(pptx_path))
    slide_count = len(prs.slides)
    stats = pptx_slide_stats(prs)
    locale = pptx_path.stem.rsplit("_", 1)[-1]
    structure_contact_sheet = pptx_text_contact_sheet(
        pptx_path, qa_dir, locale)
    render_result = render_pptx_visual_qa(
        pptx_path, expected_slides, qa_dir, reporter)
    status = slide_count == expected_slides and stats["slides_with_text"] == expected_slides and stats[
        "slides_with_notes"] == expected_slides and stats["full_slide_images"] == 0 and render_result["status"] != "FAIL"
    if status:
        reporter.pass_(f"PPTX structural QA passed: {pptx_path}")
    else:
        reporter.error(f"PPTX structural QA failed: {pptx_path}")
    return {
        "path": str(pptx_path),
        "slides": slide_count,
        "slides_with_text": stats["slides_with_text"],
        "slides_with_notes": stats["slides_with_notes"],
        "full_slide_images": stats["full_slide_images"],
        "contact_sheet": str(structure_contact_sheet),
        "visual_rendered": render_result["visual_rendered"],
        "visual_render_note": render_result["visual_render_note"],
        "visual_pdf": render_result.get("visual_pdf"),
        "visual_contact_sheet": render_result.get("visual_contact_sheet"),
        "visual_blank_pages": render_result.get("visual_blank_pages", []),
        "status": "PASS" if status else "FAIL",
    }


def render_pptx_visual_qa(pptx_path: Path, expected_slides: int, qa_dir: Path, reporter: Reporter) -> dict:
    soffice = shutil.which("soffice") or shutil.which("libreoffice")
    if not soffice:
        return {
            "visual_rendered": False,
            "visual_render_note": "LibreOffice/soffice is not available in PATH, so PPTX QA is structural plus text contact sheet.",
            "status": "SKIP",
        }
    render_root = qa_dir / "pptx-rendered"
    render_root.mkdir(parents=True, exist_ok=True)
    with tempfile.TemporaryDirectory(prefix="pptx_visual_qa_") as tmp:
        tmp_dir = Path(tmp)
        profile_dir = tmp_dir / "lo-profile"
        profile_uri = profile_dir.resolve().as_uri()
        run_checked([
            soffice,
            "--headless",
            f"-env:UserInstallation={profile_uri}",
            "--convert-to",
            "pdf",
            "--outdir",
            str(tmp_dir),
            str(pptx_path),
        ], reporter, timeout=SOFFICE_TIMEOUT_SECONDS)
        converted = tmp_dir / f"{pptx_path.stem}.pdf"
        if not converted.exists():
            raise RuntimeError(f"soffice did not produce PDF for {pptx_path}")
        visual_pdf = render_root / converted.name
        shutil.copy2(converted, visual_pdf)
    page_dir = qa_dir / "pages" / f"{pptx_path.stem}_pptx_render"
    if page_dir.exists():
        shutil.rmtree(page_dir)
    images = render_pdf_to_images(visual_pdf, page_dir, "page", reporter)
    nonwhite = [image_nonwhite_ratio(path) for path in images]
    blank_pages = [index for index, ratio in enumerate(
        nonwhite, start=1) if ratio < 0.005]
    contact_sheet = qa_dir / f"{pptx_path.stem}_visual_contact_sheet.jpg"
    make_contact_sheet(images, contact_sheet,
                       f"PPTX rendered visual QA: {pptx_path.name}")
    status = len(images) == expected_slides and not blank_pages
    if status:
        reporter.pass_(f"PPTX visual QA passed: {pptx_path}")
    else:
        reporter.error(f"PPTX visual QA failed: {pptx_path}")
    return {
        "visual_rendered": True,
        "visual_render_note": "PPTX rendered to PDF with LibreOffice/soffice and contact sheet generated.",
        "visual_pdf": str(visual_pdf),
        "visual_contact_sheet": str(contact_sheet),
        "visual_blank_pages": blank_pages,
        "status": "PASS" if status else "FAIL",
    }


def visual_qa_for_package(package_dir: Path, reporter: Reporter) -> dict:
    require_python_module("PIL", "Install Pillow to create contact sheets.")
    require_python_module("pypdf", "Install pypdf to inspect PDFs.")
    require_python_module("pptx", "Install python-pptx to inspect PPTX files.")
    html_path = customer_deck_html(package_dir)
    expected_slides = deck_slide_count(html_path)
    deck_dir = html_path.parent
    qa_dir = deck_dir / "qa"
    qa_dir.mkdir(parents=True, exist_ok=True)
    pdfs = sorted((deck_dir / "pdf").glob("**/*.pdf"))
    pptx_files = sorted((deck_dir / "pptx").glob("**/*.pptx"))
    pdf_results = [qa_pdf(path, expected_slides, qa_dir, reporter)
                   for path in pdfs]
    pptx_results = [qa_pptx(path, expected_slides, qa_dir, reporter)
                    for path in pptx_files]
    package_status = all(
        item["status"] == "PASS" for item in pdf_results + pptx_results)
    report = {
        "schema_version": "1.0.0",
        "generated_at_utc": datetime.now(timezone.utc).replace(microsecond=0).isoformat(),
        "package": str(package_dir),
        "source_html": str(html_path),
        "expected_slides": expected_slides,
        "pdf_results": pdf_results,
        "pptx_results": pptx_results,
        "status": "PASS" if package_status else "FAIL",
    }
    report_path = qa_dir / "visual-qa-report.json"
    report_path.write_text(json.dumps(
        report, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")
    records = [artifact_record("visual_qa", "qa", html_path, report_path)]
    for result in pdf_results + pptx_results:
        contact = Path(result["contact_sheet"])
        if contact.exists():
            records.append(artifact_record(
                "visual_qa", "qa", html_path, contact))
        visual_contact = result.get("visual_contact_sheet")
        if visual_contact and Path(visual_contact).exists():
            records.append(artifact_record(
                "visual_qa", "qa", html_path, Path(visual_contact)))
        visual_pdf = result.get("visual_pdf")
        if visual_pdf and Path(visual_pdf).exists():
            records.append(artifact_record(
                "visual_qa", "qa", html_path, Path(visual_pdf)))
    update_manifest_artifacts(package_dir, records, reporter)
    return report


def write_client_package(
    data: dict,
    source_path: Path,
    reference_root: Path,
    output_root: Path,
    tpid: int,
    reporter: Reporter,
) -> Path:
    package_dir = output_root / package_name(data, tpid)
    for relative_path in PACKAGE_DIRS:
        (package_dir / relative_path).mkdir(parents=True, exist_ok=True)

    snapshot = client_snapshot(data, source_path, tpid)
    snapshot_path = package_dir / DB_FILENAME
    snapshot_path.write_text(json.dumps(
        snapshot, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")
    reporter.pass_(f"wrote per-client db.json snapshot: {snapshot_path}")

    source_dir = generic_pdf_source_dir(reference_root)
    target_dir = package_dir / "EXTERNAL_Customer_Ready_Documents" / \
        "customer-ready-playbooks"
    for name in GENERIC_PDFS:
        source_file = source_dir / name
        target_file = target_dir / name
        if not source_file.exists():
            reporter.error(f"missing generic PDF source: {source_file}")
            continue
        shutil.copy2(source_file, target_file)
        reporter.pass_(f"copied generic PDF: {target_file}")

    return package_dir


def init_client(args: argparse.Namespace) -> int:
    data = load_json(args.db)
    reporter = Reporter()
    gaps = source_gaps(data, args.tpid)
    if gaps:
        reporter.source_gap(
            f"TPID {args.tpid} missing required source sections: {', '.join(gaps)}")
        if not args.allow_source_gaps:
            return reporter.exit_code()

    package_dir = write_client_package(
        data, args.db, args.reference_root, args.output_root, args.tpid, reporter)

    print(f"PACKAGE {package_dir}")
    return reporter.exit_code(args.allow_source_gaps)


def init_all_clients(args: argparse.Namespace) -> int:
    data = load_json(args.db)
    reporter = Reporter()
    created = 0
    skipped = 0

    ranking_rows = sorted(data.get("ranking_brasil", []), key=lambda row: (
        row.get("rank", 999999), str(row.get("top_parent", ""))))
    for ranking in ranking_rows:
        tpid = tpid_key(ranking)
        gaps = source_gaps(data, tpid)
        if gaps:
            skipped += 1
            reporter.source_gap(
                f"TPID {tpid} ({ranking.get('top_parent', 'unknown')}) skipped, missing required source sections: {', '.join(gaps)}")
            continue
        package_dir = write_client_package(
            data, args.db, args.reference_root, args.output_root, tpid, reporter)
        created += 1
        print(f"PACKAGE {package_dir}")

    reporter.pass_(f"initialized {created} generation-ready client packages")
    if skipped:
        reporter.source_gap(
            f"skipped {skipped} clients with missing source data")
    return reporter.exit_code(args.allow_source_gaps)


def import_transition_plan(args: argparse.Namespace) -> int:
    data = load_json(args.db)
    reporter = Reporter()
    imported = 0
    skipped = 0

    for source_dir in transition_source_folders(args.source, args.client_folder):
        try:
            tpid = resolve_transition_tpid(source_dir.name)
        except ValueError as exc:
            skipped += 1
            reporter.error(str(exc))
            continue

        gaps = source_gaps(data, tpid)
        if gaps:
            skipped += 1
            reporter.source_gap(
                f"folder {source_dir.name} TPID {tpid} skipped, missing required source sections: {', '.join(gaps)}"
            )
            if not args.allow_source_gaps:
                continue

        package_dir = write_client_package(
            data, args.db, args.reference_root, args.output_root, tpid, reporter)
        custom_records = copy_custom_deliverables(
            source_dir, package_dir, reporter)
        if len(custom_records) != len(SOURCE_CUSTOM_DELIVERABLES):
            skipped += 1
            continue
        artifacts = custom_records + \
            generic_artifact_records(args.reference_root, package_dir)
        write_manifest(data, args.db, source_dir,
                       package_dir, tpid, artifacts, reporter)
        imported += 1
        print(f"PACKAGE {package_dir}")

    reporter.pass_(f"imported {imported} transition-plan packages")
    if skipped:
        reporter.warn(f"skipped {skipped} transition-plan folders")
    return reporter.exit_code(args.allow_source_gaps)


def load_package_tpid(package_dir: Path) -> int:
    package_db = load_json(package_dir / DB_FILENAME)
    client = package_db.get("client") or {}
    if "tpid" not in client:
        raise ValueError(
            f"Package db.json missing client.tpid: {package_dir / DB_FILENAME}")
    return int(client["tpid"])


def validate_snapshot(source_data: dict, source_path: Path, package_dir: Path, reporter: Reporter) -> None:
    package_db_path = package_dir / DB_FILENAME
    if not package_db_path.exists():
        reporter.error(f"package missing db.json: {package_db_path}")
        return
    package_db = load_json(package_db_path)
    tpid = int(package_db.get("client", {}).get("tpid", 0))
    expected = client_snapshot(source_data, source_path, tpid)
    if "source" in package_db and "source" in expected:
        expected["source"]["file"] = package_db["source"].get("file")
    if package_db != expected:
        reporter.error(
            f"package db.json does not match source db.json for TPID {tpid}")
    else:
        reporter.pass_(
            f"package db.json matches source db.json for TPID {tpid}")


def non_archive_matches(directory: Path, pattern: re.Pattern[str]) -> list[Path]:
    if not directory.exists():
        return []
    matches = []
    for path in directory.iterdir():
        if path.is_file() and pattern.match(path.name):
            matches.append(path)
    return sorted(matches)


def validate_package_dirs(package_dir: Path, reporter: Reporter) -> None:
    for relative_path in PACKAGE_DIRS:
        path = package_dir / relative_path
        if path.is_dir():
            reporter.pass_(f"required directory exists: {path}")
        else:
            reporter.error(f"required directory missing: {path}")


def validate_generic_pdfs(package_dir: Path, reporter: Reporter) -> None:
    playbooks_dir = package_dir / \
        "EXTERNAL_Customer_Ready_Documents" / "customer-ready-playbooks"
    for name in GENERIC_PDFS:
        path = playbooks_dir / name
        if path.exists():
            reporter.pass_(f"generic PDF present: {path.name}")
        else:
            reporter.error(f"generic PDF missing: {path}")


def validate_custom_deliverables(package_dir: Path, reporter: Reporter) -> None:
    for label, (relative_dir, pattern) in CUSTOM_DELIVERABLES.items():
        matches = non_archive_matches(package_dir / relative_dir, pattern)
        if len(matches) == 1:
            reporter.pass_(
                f"{label} latest deliverable present: {matches[0].name}")
        elif len(matches) == 0:
            reporter.error(
                f"{label} deliverable missing in {package_dir / relative_dir}")
        else:
            names = ", ".join(path.name for path in matches)
            reporter.error(
                f"{label} has multiple current versions outside archive: {names}")


def validate_manifest(package_dir: Path, reporter: Reporter) -> None:
    manifest_path = package_dir / MANIFEST_FILENAME
    if not manifest_path.exists():
        reporter.error(f"package missing manifest: {manifest_path}")
        return
    manifest = load_json(manifest_path)
    artifacts = manifest.get("artifacts")
    if not isinstance(artifacts, list):
        reporter.error(f"manifest artifacts must be a list: {manifest_path}")
        return
    for artifact in artifacts:
        destination = Path(str(artifact.get("destination_path", "")))
        expected_sha = artifact.get("sha256")
        if not destination.exists():
            reporter.error(f"manifest artifact missing: {destination}")
            continue
        actual_sha = sha256_file(destination)
        if actual_sha != expected_sha:
            reporter.error(f"manifest checksum mismatch for {destination}")
        else:
            reporter.pass_(f"manifest checksum ok: {destination.name}")


def validate_package(args: argparse.Namespace) -> int:
    data = load_json(args.db)
    package_dir = args.package
    reporter = Reporter()

    if not package_dir.exists():
        reporter.error(f"package directory not found: {package_dir}")
        return reporter.exit_code()

    validate_package_dirs(package_dir, reporter)

    validate_snapshot(data, args.db, package_dir, reporter)

    try:
        tpid = load_package_tpid(package_dir)
    except ValueError as exc:
        reporter.error(str(exc))
        return reporter.exit_code()

    gaps = source_gaps(data, tpid)
    if gaps:
        reporter.source_gap(
            f"TPID {tpid} missing required source sections: {', '.join(gaps)}")

    validate_generic_pdfs(package_dir, reporter)

    if not args.scaffold_only:
        validate_custom_deliverables(package_dir, reporter)
        if not args.no_manifest:
            validate_manifest(package_dir, reporter)

    return reporter.exit_code(args.allow_source_gaps)


def validate_packages(args: argparse.Namespace) -> int:
    reporter = Reporter()
    if not args.output_root.exists():
        reporter.error(f"output root not found: {args.output_root}")
        return reporter.exit_code()
    packages = [path for path in args.output_root.iterdir() if path.is_dir()]
    if not packages:
        reporter.warn(f"no client packages found under {args.output_root}")
        return 0
    exit_code = 0
    for package in sorted(packages):
        package_args = argparse.Namespace(
            db=args.db,
            reference_root=args.reference_root,
            package=package,
            scaffold_only=args.scaffold_only,
            allow_source_gaps=args.allow_source_gaps,
            no_manifest=args.no_manifest,
        )
        result = validate_package(package_args)
        exit_code = max(exit_code, result)
    return exit_code


def generate_deck_derivatives(args: argparse.Namespace) -> int:
    reporter = Reporter()
    if args.package:
        packages = args.package
    else:
        if not args.output_root.exists():
            reporter.error(f"output root not found: {args.output_root}")
            return reporter.exit_code()
        packages = sorted(
            path for path in args.output_root.iterdir() if path.is_dir())
    pdf_locales = tuple(args.pdf_locale or DEFAULT_PDF_LOCALES)
    pptx_locales = tuple(args.pptx_locale or DECK_DERIVATIVE_LOCALES)
    processed = 0
    for package_dir in packages:
        try:
            generate_deck_derivatives_for_package(
                package_dir, pdf_locales, pptx_locales, reporter)
            processed += 1
        except (ValueError, RuntimeError, FileNotFoundError) as exc:
            reporter.error(f"{package_dir}: {exc}")
            if not args.keep_going:
                return reporter.exit_code()
    reporter.pass_(f"generated deck derivatives for {processed} packages")
    return reporter.exit_code()


def visual_qa_deck_derivatives(args: argparse.Namespace) -> int:
    reporter = Reporter()
    if args.package:
        packages = args.package
    else:
        if not args.output_root.exists():
            reporter.error(f"output root not found: {args.output_root}")
            return reporter.exit_code()
        packages = sorted(
            path for path in args.output_root.iterdir() if path.is_dir())
    reports = []
    for package_dir in packages:
        try:
            reports.append(visual_qa_for_package(package_dir, reporter))
        except (RuntimeError, ValueError, FileNotFoundError) as exc:
            reporter.error(f"{package_dir}: {exc}")
            if not args.keep_going:
                return reporter.exit_code()
    summary = {
        "schema_version": "1.0.0",
        "generated_at_utc": datetime.now(timezone.utc).replace(microsecond=0).isoformat(),
        "packages": len(reports),
        "passed": sum(1 for report in reports if report.get("status") == "PASS"),
        "failed": sum(1 for report in reports if report.get("status") != "PASS"),
        "reports": reports,
    }
    summary_path = args.output_root / "visual-qa-summary.json"
    summary_path.write_text(json.dumps(
        summary, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")
    if summary["failed"]:
        reporter.error(f"visual QA failed for {summary['failed']} packages")
    else:
        reporter.pass_(f"visual QA passed for {summary['passed']} packages")
    reporter.pass_(f"wrote visual QA summary: {summary_path}")
    return reporter.exit_code()


def command_validate_db(args: argparse.Namespace) -> int:
    reporter = Reporter()
    validate_db(load_json(args.db), reporter)
    return reporter.exit_code(args.allow_source_gaps)


def build_parser() -> argparse.ArgumentParser:
    parser = argparse.ArgumentParser(description=__doc__)
    sub = parser.add_subparsers(dest="command", required=True)

    validate_db_parser = sub.add_parser(
        "validate-db", help="Validate canonical db.json coverage and reconciliation.")
    validate_db_parser.add_argument("--db", type=Path, required=True)
    validate_db_parser.add_argument("--allow-source-gaps", action="store_true")
    validate_db_parser.set_defaults(func=command_validate_db)

    init_parser = sub.add_parser(
        "init-client", help="Create a client package folder and source snapshot.")
    init_parser.add_argument("--db", type=Path, required=True)
    init_parser.add_argument(
        "--reference-root", type=Path, default=DEFAULT_REFERENCE_ROOT)
    init_parser.add_argument("--output-root", type=Path,
                             default=DEFAULT_OUTPUT_ROOT)
    init_parser.add_argument("--tpid", type=int, required=True)
    init_parser.add_argument("--allow-source-gaps", action="store_true")
    init_parser.set_defaults(func=init_client)

    init_all_parser = sub.add_parser(
        "init-all-clients", help="Create package folders and source snapshots for every generation-ready ranked client."
    )
    init_all_parser.add_argument("--db", type=Path, required=True)
    init_all_parser.add_argument(
        "--reference-root", type=Path, default=DEFAULT_REFERENCE_ROOT)
    init_all_parser.add_argument("--output-root", type=Path,
                                 default=DEFAULT_OUTPUT_ROOT)
    init_all_parser.add_argument("--allow-source-gaps", action="store_true")
    init_all_parser.set_defaults(func=init_all_clients)

    import_parser = sub.add_parser(
        "import-transition-plan", help="Import the audited 00-Customers_Transition_Plan materials into package folders."
    )
    import_parser.add_argument("--db", type=Path, required=True)
    import_parser.add_argument(
        "--reference-root", type=Path, default=DEFAULT_REFERENCE_ROOT)
    import_parser.add_argument("--source", type=Path,
                               default=DEFAULT_TRANSITION_PLAN_SOURCE)
    import_parser.add_argument("--output-root", type=Path,
                               default=DEFAULT_OUTPUT_ROOT)
    import_parser.add_argument("--client-folder", action="append")
    import_parser.add_argument("--allow-source-gaps", action="store_true")
    import_parser.set_defaults(func=import_transition_plan)

    package_parser = sub.add_parser(
        "validate-package", help="Validate one completed client package.")
    package_parser.add_argument("--db", type=Path, required=True)
    package_parser.add_argument(
        "--reference-root", type=Path, default=DEFAULT_REFERENCE_ROOT)
    package_parser.add_argument("--package", type=Path, required=True)
    package_parser.add_argument("--scaffold-only", action="store_true")
    package_parser.add_argument("--no-manifest", action="store_true")
    package_parser.add_argument("--allow-source-gaps", action="store_true")
    package_parser.set_defaults(func=validate_package)

    packages_parser = sub.add_parser(
        "validate-packages", help="Validate all client packages under the output root.")
    packages_parser.add_argument("--db", type=Path, required=True)
    packages_parser.add_argument(
        "--reference-root", type=Path, default=DEFAULT_REFERENCE_ROOT)
    packages_parser.add_argument(
        "--output-root", type=Path, default=DEFAULT_OUTPUT_ROOT)
    packages_parser.add_argument("--scaffold-only", action="store_true")
    packages_parser.add_argument("--no-manifest", action="store_true")
    packages_parser.add_argument("--allow-source-gaps", action="store_true")
    packages_parser.set_defaults(func=validate_packages)

    derivatives_parser = sub.add_parser(
        "generate-deck-derivatives", help="Generate PDF and native PPTX derivatives for customer deck HTML files."
    )
    derivatives_parser.add_argument(
        "--output-root", type=Path, default=DEFAULT_OUTPUT_ROOT)
    derivatives_parser.add_argument("--package", type=Path, action="append")
    derivatives_parser.add_argument(
        "--pdf-locale", choices=DECK_DERIVATIVE_LOCALES, action="append")
    derivatives_parser.add_argument(
        "--pptx-locale", choices=DECK_DERIVATIVE_LOCALES, action="append")
    derivatives_parser.add_argument("--keep-going", action="store_true")
    derivatives_parser.set_defaults(func=generate_deck_derivatives)

    visual_qa_parser = sub.add_parser(
        "visual-qa-deck-derivatives", help="Render package deck PDFs, inspect PPTX files, and create visual QA contact sheets."
    )
    visual_qa_parser.add_argument(
        "--output-root", type=Path, default=DEFAULT_OUTPUT_ROOT)
    visual_qa_parser.add_argument("--package", type=Path, action="append")
    visual_qa_parser.add_argument("--keep-going", action="store_true")
    visual_qa_parser.set_defaults(func=visual_qa_deck_derivatives)

    return parser


def main(argv: list[str] | None = None) -> int:
    parser = build_parser()
    args = parser.parse_args(argv)
    return args.func(args)


if __name__ == "__main__":
    sys.exit(main())
