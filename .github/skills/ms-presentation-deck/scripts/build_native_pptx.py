#!/usr/bin/env python3
"""Build native editable PPTX files from an ms-presentation-deck HTML source.

This is intentionally conservative: it reads the HTML deck, resolves visible
`data-i18n` content for one locale, and writes native PowerPoint text boxes,
cards, and speaker notes. It never screenshots a slide.
"""

from __future__ import annotations

import argparse
import json
import re
import subprocess
import sys
from html import unescape
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
    from pptx.util import Emu, Pt
except ImportError as exc:
    sys.exit(
        f"ERROR: missing dependency ({exc.name}).\n"
        "  pip install -r .github/skills/ms-presentation-deck/scripts/requirements.txt"
    )

EMU_IN = 914400
PW, PH = 13.333, 7.5
LEFT, RIGHT = 0.62, 12.71
TOP, BOTTOM = 0.70, 6.95

COLORS = {
    "red": RGBColor(0xF2, 0x50, 0x22),
    "green": RGBColor(0x7F, 0xBA, 0x00),
    "blue": RGBColor(0x00, 0xA4, 0xEF),
    "yellow": RGBColor(0xFF, 0xB9, 0x00),
    "ink": RGBColor(0x1A, 0x1A, 0x1A),
    "ink2": RGBColor(0x4A, 0x4A, 0x4A),
    "ink3": RGBColor(0x82, 0x80, 0x7A),
    "paper": RGBColor(0xFF, 0xFF, 0xFF),
    "soft": RGBColor(0xF7, 0xF7, 0xF5),
    "rule": RGBColor(0xE5, 0xE5, 0xE0),
    "dark": RGBColor(0x14, 0x14, 0x14),
    "dark2": RGBColor(0x2E, 0x2E, 0x2A),
    "white": RGBColor(0xFF, 0xFF, 0xFF),
}

LOCALES = ("pt-BR", "en", "es")


def inch(value: float) -> Emu:
    return Emu(int(value * EMU_IN))


def fail(message: str) -> None:
    print(f"ERROR: {message}", file=sys.stderr)
    raise SystemExit(1)


def strip_html(value: str) -> str:
    value = re.sub(r"<br\s*/?>", "\n", value or "", flags=re.I)
    value = re.sub(r"</(p|div|li|h\d)>", "\n", value, flags=re.I)
    value = re.sub(r"<[^>]+>", " ", value)
    value = unescape(value)
    value = value.replace("·", " · ")
    value = re.sub(r"[ \t]+", " ", value)
    value = re.sub(r"\n\s+", "\n", value)
    value = re.sub(r"\n{3,}", "\n\n", value)
    return value.strip()


def load_i18n(html: str) -> dict:
    marker = "const I18N = "
    if marker not in html:
        fail("HTML source does not contain `const I18N = ...`")
    start = html.index(marker) + len(marker)
    try:
        data, _ = json.JSONDecoder().raw_decode(html[start:])
    except json.JSONDecodeError as exc:
        fail(f"could not parse I18N JSON: {exc}")
    return data


def deck_base(path: Path) -> str:
    return path.stem[:-6] if path.stem.endswith("_multi") else path.stem


def sections(html: str) -> list[str]:
    found = re.findall(
        r"<section\b(?=[^>]*\bclass=[\"'][^\"']*\bslide\b).*?</section>",
        html,
        flags=re.I | re.S,
    )
    if not found:
        fail("HTML source has no slide sections")
    return found


def resolve(data: dict, locale: str, key: str):
    current = data.get(locale, {})
    for part in key.split("."):
        current = current.get(part, "") if isinstance(current, dict) else ""
    return current


def visible_keys(section: str) -> list[str]:
    keys = re.findall(r"data-i18n=[\"']([^\"']+)[\"']", section)
    keys += re.findall(r"data-i18n-list=[\"']([^\"']+)[\"']", section)
    out: list[str] = []
    for key in keys:
        if key not in out:
            out.append(key)
    return out


def validate_locale_coverage(html: str, data: dict, locale: str) -> None:
    keys = sorted({key for section in sections(html) for key in visible_keys(section)})
    if locale != "en" and not keys:
        fail(f"{locale} cannot be exported: HTML has no visible data-i18n keys")
    missing = [key for key in keys if resolve(data, "en", key) not in (None, "") and resolve(data, locale, key) in (None, "")]
    if missing:
        fail(f"{locale} cannot be exported: visible translations missing for keys {missing[:12]}")


def note_keys(html: str, total: int) -> list[str]:
    match = re.search(r"const NOTE_MAP = (\[[^\]]*\]);", html, flags=re.S)
    if not match:
        return [f"s{i}" for i in range(1, total + 1)]
    note_map = json.loads(match.group(1))
    return [note_map[i - 1] or f"s{i}" if i - 1 < len(note_map) else f"s{i}" for i in range(1, total + 1)]


def slide_values(section: str, data: dict, locale: str) -> list[str]:
    values: list[str] = []
    for key in visible_keys(section):
        value = resolve(data, locale, key)
        if isinstance(value, list):
            for item in value:
                text = strip_html(str(item))
                if text and text not in values:
                    values.append(text)
            continue
        text = strip_html(str(value))
        if text and text not in values:
            values.append(text)
    if values:
        return values
    # Fallback for older decks with static visible text.
    raw = re.sub(r"<script.*?</script>|<style.*?</style>", " ", section, flags=re.I | re.S)
    text = strip_html(raw)
    lines = [line.strip() for line in text.splitlines() if len(line.strip()) > 1]
    return lines[:12]


def clean_note(note: str) -> str:
    return re.sub(r"\s+", " ", (note or "").replace("**", "").replace("*", "")).strip()


def notes_to_values(note: str, locale: str, page: int) -> list[str]:
    text = clean_note(note)
    text = re.sub(r"\[[^\]]+\]", " ", text)
    text = re.sub(r"\s+", " ", text).strip()
    if not text:
        return [f"Slide {page}"]
    sentences = re.split(r"(?<=[.!?])\s+", text)
    sentences = [sentence.strip() for sentence in sentences if len(sentence.strip()) > 12]
    labels = {"pt-BR": "Notas do apresentador", "en": "Speaker notes", "es": "Notas del presentador"}
    title = sentences[0] if sentences else text[:120]
    body = sentences[1:7] if len(sentences) > 1 else [text]
    return [labels.get(locale, "Speaker notes"), title, *body]


def add_box(slide, x, y, w, h, fill=None, line=None):
    shape = slide.shapes.add_shape(1, inch(x), inch(y), inch(w), inch(h))
    if fill is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(0.8)
    return shape


def add_text(slide, text, x, y, w, h, size=16, color=None, bold=False, font="Segoe UI", align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(inch(x), inch(y), inch(w), inch(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = None
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text or ""
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.name = font
    run.font.color.rgb = color or COLORS["ink"]
    return box


def brand(slide, page: int, total: int, dark: bool = False) -> None:
    y = 0.28
    x = 0.28
    size = 0.11
    gap = 0.02
    add_box(slide, x, y, size, size, COLORS["red"])
    add_box(slide, x + size + gap, y, size, size, COLORS["green"])
    add_box(slide, x, y + size + gap, size, size, COLORS["blue"])
    add_box(slide, x + size + gap, y + size + gap, size, size, COLORS["yellow"])
    color = COLORS["white"] if dark else COLORS["ink"]
    add_text(slide, "PAULA SILVA  |  SOFTWARE GLOBAL BLACK BELT", 0.62, 0.28, 5.0, 0.25, 7.5, color, False, "Consolas")
    add_text(slide, f"{page:02d} / {total}", 11.6, 7.08, 1.1, 0.22, 8, color, False, "Consolas", PP_ALIGN.RIGHT)


def add_notes(slide, note: str) -> None:
    slide.notes_slide.notes_text_frame.text = clean_note(note)


def slide_copy(values: list[str], page: int) -> tuple[str, str, list[str]]:
    if not values:
        return "", f"Slide {page}", []
    eyebrow = values[0]
    title = values[1] if len(values) > 1 else values[0]
    body = values[2:] if len(values) > 2 else []
    if len(title) < 18 and body:
        title = body.pop(0)
    body = [item for item in body if item and item != title][:8]
    return eyebrow, title, body or [title]


def add_body_items(slide, items: list[str], y: float, accent, body_color) -> None:
    for item in items:
        if y > 6.45:
            break
        add_box(slide, LEFT, y, 0.13, 0.13, accent)
        size = 13.5 if len(item) < 170 else 11.5
        height = 0.58 if len(item) < 170 else 0.82
        add_text(slide, item[:420], LEFT + 0.28, y - 0.04, 11.2, height, size, body_color)
        y += 0.70 if len(item) < 170 else 0.92


def render_slide(prs: Presentation, values: list[str], note: str, page: int, total: int, dark: bool) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = COLORS["dark"] if dark else COLORS["paper"]
    add_box(slide, 0, 0, PW, PH, bg)
    brand(slide, page, total, dark)
    title_color = COLORS["white"] if dark else COLORS["ink"]
    body_color = COLORS["white"] if dark else COLORS["ink2"]
    accent = COLORS["yellow"] if dark else COLORS["blue"]

    eyebrow, title, body = slide_copy(values, page)

    add_text(slide, eyebrow.upper()[:110], LEFT, 0.95, 11.5, 0.32, 9.5, accent, True, "Consolas")
    add_text(slide, title[:220], LEFT, 1.35, 11.8, 1.25, 28 if len(title) < 95 else 23, title_color, False)
    add_body_items(slide, body, 2.75, accent, body_color)
    add_notes(slide, note)


def build(html_path: Path, locale: str, output: Path, allow_notes_derived_content: bool = False) -> Path:
    html = html_path.read_text(encoding="utf-8")
    data = load_i18n(html)
    if locale not in data:
        fail(f"locale {locale} missing in {html_path}")
    if not allow_notes_derived_content:
        validate_locale_coverage(html, data, locale)
    slide_sections = sections(html)
    keys = note_keys(html, len(slide_sections))
    notes = data[locale].get("notes", {})

    prs = Presentation()
    prs.slide_width = inch(PW)
    prs.slide_height = inch(PH)
    total = len(slide_sections)
    for page, section in enumerate(slide_sections, start=1):
        note = notes.get(keys[page - 1])
        if not note:
            fail(f"missing {locale} note for slide {page}")
        dark = "slide--dark" in section or "slide--divider" in section
        values = notes_to_values(note, locale, page) if allow_notes_derived_content else slide_values(section, data, locale)
        render_slide(prs, values, note, page, total, dark)

    output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output)
    if not output.is_file() or output.stat().st_size == 0:
        fail(f"PPTX was not written: {output}")
    return output


def validate(html_path: Path, locale: str, output: Path, allow_notes_derived_content: bool = False) -> None:
    cmd = [
        "python3",
        ".github/skills/ms-presentation-deck/scripts/validate_derivatives.py",
        "--html",
        str(html_path),
        "--locale",
        locale,
        "--pptx",
        str(output),
    ]
    if allow_notes_derived_content:
        cmd.append("--allow-notes-derived-content")
    result = subprocess.run(cmd, cwd=Path.cwd(), capture_output=True, text=True)
    if result.returncode != 0:
        fail((result.stdout + result.stderr).strip())


def main() -> int:
    parser = argparse.ArgumentParser(description="Build native editable PPTX derivatives from an HTML deck")
    parser.add_argument("--html", required=True, help="Source HTML deck")
    parser.add_argument("--locale", choices=LOCALES, help="Locale to export")
    parser.add_argument("--all-locales", action="store_true", help="Build pt-BR, en, and es")
    parser.add_argument(
        "--allow-notes-derived-content",
        action="store_true",
        help="Use localized speaker notes as visible slide content when the HTML visible text is not localized",
    )
    parser.add_argument("--output", help="Output PPTX path for single-locale export")
    args = parser.parse_args()
    if not args.locale and not args.all_locales:
        fail("provide --locale or --all-locales")
    html_path = Path(args.html)
    base = deck_base(html_path)
    locales = LOCALES if args.all_locales else (args.locale,)
    for locale in locales:
        output = Path(args.output) if args.output else html_path.parent / "pptx" / base / f"{base}_{locale}.pptx"
        written = build(html_path, locale, output, args.allow_notes_derived_content)
        validate(html_path, locale, written, args.allow_notes_derived_content)
        print(f"OK: wrote {written}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())